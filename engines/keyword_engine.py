"""
engines/keyword_engine.py
Keyword search engine — extracts text from images (OCR) and documents,
then searches for user-provided keywords.

Supported file types:
  Images  — via EasyOCR         (pip install easyocr)
  PDF     — via pdfplumber       (pip install pdfplumber)
  Word    — via python-docx      (pip install python-docx)
  PPT     — via python-pptx      (pip install python-pptx)
  TXT/CSV/HTML/JSON — built-in, no install needed
"""

import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Callable

from config import IMG_EXTENSIONS, DOC_EXTENSIONS, get_logger
from engines.core_engine import ForensicEngine

logger = get_logger("KeywordEngine")

# Optional library imports
easyocr_ok    = False
pdfplumber_ok = False
docx_ok       = False
pptx_ok       = False

try:
    import easyocr as _easyocr
    easyocr_ok = True
except ImportError:
    pass

try:
    import pdfplumber as _pdfplumber
    pdfplumber_ok = True
except ImportError:
    pass

try:
    from docx import Document as _DocxDocument
    docx_ok = True
except ImportError:
    pass

try:
    from pptx import Presentation as _PptxPresentation
    pptx_ok = True
except ImportError:
    pass


class KeywordSearchEngine:
    """Search keywords across images (via OCR) and documents."""

    def __init__(self):
        self.ocr_available  = easyocr_ok
        self.pdf_available  = pdfplumber_ok
        self.docx_available = docx_ok
        self.pptx_available = pptx_ok
        self._ocr_reader    = None

    @property
    def available(self) -> bool:
        # Always True — plain text files need zero extra libraries
        return True

    # ── OCR loader ────────────────────────────────────────────────
    def load_ocr(self, log_fn: Optional[Callable] = None) -> bool:
        """Lazy-load EasyOCR. Call once before scanning images."""
        if self._ocr_reader is not None:
            return True
        if not self.ocr_available:
            return False
        try:
            if log_fn:
                log_fn("Loading EasyOCR model (first run ~200 MB)…", "info")
            self._ocr_reader = _easyocr.Reader(
                ["en"], gpu=False, verbose=False)
            if log_fn:
                log_fn("✔ EasyOCR ready.", "match")
            return True
        except Exception as e:
            logger.error(f"EasyOCR load failed: {e}")
            self.ocr_available = False
            return False

    # ── Text extractors ───────────────────────────────────────────
    def extract_image_text(self, filepath: str) -> str:
        if self._ocr_reader is None:
            return ""
        try:
            parts = self._ocr_reader.readtext(filepath, detail=0, paragraph=True)
            return " ".join(parts)
        except Exception as e:
            logger.debug(f"OCR error {filepath}: {e}")
            return ""

    def extract_pdf_text(self, filepath: str) -> str:
        if not self.pdf_available:
            return ""
        try:
            with _pdfplumber.open(filepath) as pdf:
                return "\n".join(
                    p.extract_text() for p in pdf.pages
                    if p.extract_text())
        except Exception as e:
            logger.debug(f"PDF extract error {filepath}: {e}")
            return ""

    def extract_docx_text(self, filepath: str) -> str:
        if not self.docx_available:
            return ""
        try:
            doc = _DocxDocument(filepath)
            return "\n".join(
                p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.debug(f"DOCX extract error {filepath}: {e}")
            return ""

    def extract_pptx_text(self, filepath: str) -> str:
        if not self.pptx_available:
            return ""
        try:
            prs = _PptxPresentation(filepath)
            return "\n".join(
                shape.text
                for slide in prs.slides
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip())
        except Exception as e:
            logger.debug(f"PPTX extract error {filepath}: {e}")
            return ""

    def extract_txt_text(self, filepath: str) -> str:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""

    def extract_text(self, filepath: str) -> str:
        """Route to the correct extractor based on file extension."""
        ext = Path(filepath).suffix.lower()
        if ext in IMG_EXTENSIONS:       return self.extract_image_text(filepath)
        elif ext == ".pdf":             return self.extract_pdf_text(filepath)
        elif ext in (".docx", ".doc"):  return self.extract_docx_text(filepath)
        elif ext in (".pptx", ".ppt"):  return self.extract_pptx_text(filepath)
        else:                           return self.extract_txt_text(filepath)

    # ── Keyword matching ──────────────────────────────────────────
    @staticmethod
    def search_keywords(text: str,
                        keywords: List[str],
                        case_sensitive: bool = False) -> Dict[str, List[str]]:
        """
        Find keywords in text. Returns {keyword: [context_sentences]}.
        Uses whole-word matching by default.
        """
        if not text or not keywords:
            return {}
        flags = 0 if case_sensitive else re.IGNORECASE
        sents = [s.strip() for s in re.split(r'[.\n!?;]', text) if s.strip()]
        hits  = {}
        for kw in keywords:
            kw = kw.strip()
            if not kw:
                continue
            pat = re.compile(r'\b' + re.escape(kw) + r'\b', flags)
            matched = [s for s in sents if pat.search(s)]
            if matched:
                hits[kw] = matched[:5]   # max 5 context sentences per keyword
        return hits

    # ── Full folder scan ──────────────────────────────────────────
    def scan_folder(self,
                    folder: str,
                    keywords: List[str],
                    do_images: bool = True,
                    do_docs: bool = True,
                    case_sensitive: bool = False,
                    progress_cb: Optional[Callable] = None,
                    log_fn: Optional[Callable] = None,
                    stop_fn: Optional[Callable] = None) -> List[Dict[str, Any]]:
        """
        Scan an entire folder for keyword matches.

        Args:
            folder         : path to the folder to scan
            keywords       : list of keywords to search for
            do_images      : include image files (requires EasyOCR)
            do_docs        : include document files
            case_sensitive : whether matching is case-sensitive
            progress_cb    : callback(done, total, filename) for UI progress
            log_fn         : callback(msg, tag) for live log output
            stop_fn        : callable that returns True when scan should stop.
                             Checked BEFORE each file so stop is near-instant.

        Returns:
            List of match dicts, each containing filepath, keywords_found,
            metadata, hashes etc.
        """
        def _log(msg, tag="info"):
            logger.info(msg)
            if log_fn:
                log_fn(f"  {msg}", tag)

        def _should_stop() -> bool:
            return stop_fn is not None and stop_fn()

        # Collect files
        folder_path = Path(folder)
        all_files   = []
        if do_images:
            for ext in IMG_EXTENSIONS:
                all_files += list(folder_path.rglob(f"*{ext}"))
                all_files += list(folder_path.rglob(f"*{ext.upper()}"))
        if do_docs:
            for ext in DOC_EXTENSIONS:
                all_files += list(folder_path.rglob(f"*{ext}"))
                all_files += list(folder_path.rglob(f"*{ext.upper()}"))

        # Case-insensitive deduplication (important on Windows)
        seen, deduped = set(), []
        for f in all_files:
            key = str(f).lower()
            if key not in seen:
                seen.add(key)
                deduped.append(f)
        all_files = deduped
        total     = len(all_files)

        _log(f"Keyword scan — {total} file(s) found "
             f"(images={do_images}, docs={do_docs})", "info")
        if total == 0:
            _log("No files found to search.", "error")
            return []

        results = []
        for i, fp in enumerate(all_files):

            # ── Stop check BEFORE processing each file ────────────
            # This is the key fix — checked before OCR/extraction begins
            # so stop is instant rather than waiting for the current file
            if _should_stop():
                _log(f"⏹ Scan stopped by user after {i} file(s).", "skip")
                break

            if progress_cb:
                progress_cb(i + 1, total, fp.name)

            ext = fp.suffix.lower()
            if ext in IMG_EXTENSIONS and not self.ocr_available:
                continue   # Skip images if OCR not installed

            try:
                text = self.extract_text(str(fp))
            except Exception as e:
                _log(f"Extract error: {fp.name} — {e}", "error")
                continue

            if not text or not text.strip():
                logger.debug(f"No text from: {fp.name}")
                continue

            hits = self.search_keywords(text, keywords, case_sensitive)
            if not hits:
                continue

            ftype = "image" if ext in IMG_EXTENSIONS else "document"
            try:
                metadata = ForensicEngine.extract_metadata(str(fp))
                hashes   = ForensicEngine.compute_hashes(str(fp))
            except Exception:
                metadata = {"filename": fp.name,
                            "folder":   str(fp.parent),
                            "filepath": str(fp)}
                hashes   = {}

            results.append({
                "filepath":       str(fp),
                "filename":       fp.name,
                "filetype":       ftype,
                "extension":      ext,
                "keywords_found": hits,
                "keyword_count":  len(hits),
                "total_hits":     sum(len(v) for v in hits.values()),
                "metadata":       metadata,
                "hashes":         hashes,
            })
            _log(f"✔ KEYWORD MATCH  [{', '.join(hits.keys())}]  →  {fp.name}",
                 "match")

        stopped = _should_stop()
        _log(
            f"Keyword scan {'stopped' if stopped else 'complete'} — "
            f"{len(results)} match(es) from {i + 1 if total else 0} file(s) checked.",
            "skip" if stopped else ("match" if results else "info")
        )
        return results
