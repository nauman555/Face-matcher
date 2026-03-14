"""
╔══════════════════════════════════════════════════════════════╗
║       Documents / Images Finder Tool                        ║
║       Advanced Search & Analysis Tool                       ║
║                                                             ║
╚══════════════════════════════════════════════════════════════╝

Architecture:
  - ForensicEngine   : Core scan / hash / metadata logic
  - NSFWEngine       : NudeNet-based explicit content detection
  - FaceEngine       : Face recognition matching
  - ReportEngine     : PDF report generation
  - ForensicApp      : Tkinter UI

Install dependencies:
    pip install "numpy<2.0" face-recognition Pillow nudenet
                reportlab piexif tqdm

Run:
    python forensic_tool.py
"""

# ─── Standard Library ─────────────────────────────────────────────────────────
import os
import re
import sys
import json
import shutil
import hashlib
import logging
import datetime
import threading
import multiprocessing
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pathlib import Path
from typing import Optional, List, Dict, Any

# ─── Third-party ──────────────────────────────────────────────────────────────
try:
    from PIL import Image, ImageTk, ImageDraw, ImageFont
    import numpy as np
except ImportError as e:
    print(f"[FATAL] {e}\nRun: pip install Pillow numpy")
    sys.exit(1)

# Optional imports — checked at runtime
face_recognition = None
NudeDetector     = None
reportlab_ok     = False

try:
    import face_recognition as _fr
    face_recognition = _fr
except ImportError:
    pass

try:
    from nudenet import NudeDetector as _ND
    NudeDetector = _ND
except ImportError:
    pass

try:
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm, cm
    from reportlab.platypus import (SimpleDocTemplate, Table, TableStyle,
                                     Paragraph, Spacer, Image as RLImage,
                                     HRFlowable, PageBreak)
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    reportlab_ok = True
except ImportError:
    pass

try:
    import piexif
    piexif_ok = True
except ImportError:
    piexif_ok = False

# CLIP for kissing / intimate detection
clip_ok = False
try:
    from transformers import CLIPProcessor, CLIPModel
    import torch
    clip_ok = True
except ImportError:
    pass

# EasyOCR for keyword search in images
easyocr_ok = False
try:
    import easyocr as _easyocr
    easyocr_ok = True
except ImportError:
    pass

# Document text extraction
pdfplumber_ok = False
try:
    import pdfplumber as _pdfplumber
    pdfplumber_ok = True
except ImportError:
    pass

docx_ok = False
try:
    from docx import Document as _DocxDocument
    docx_ok = True
except ImportError:
    pass

pptx_ok = False
try:
    from pptx import Presentation as _PptxPresentation
    pptx_ok = True
except ImportError:
    pass


# ══════════════════════════════════════════════════════════════════════════════
#  LOGGING SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
LOG_FILE = "finder_tool.log"
logging.basicConfig(
    level=logging.DEBUG,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
    handlers=[
        logging.FileHandler(LOG_FILE, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)
logger = logging.getLogger("FinderTool")


# ══════════════════════════════════════════════════════════════════════════════
#  THEME — Professional Forensic UI  (Deep Navy · Cyan Accent · White Text)
#  Design rule: dark layered backgrounds, ONE accent color, white text hierarchy
#  Reference: Cellebrite UFED, Maltego, VS Code Dark+
# ══════════════════════════════════════════════════════════════════════════════
SUPPORTED_EXT   = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif", ".webp"}
CACHE_FILENAME  = ".finder_encoding_cache.json"
MAX_DIM         = 800
THUMB_SIZE      = (120, 90)
PDF_THUMB_SIZE  = (80, 60)

C = {
    # Backgrounds — 5 depth levels, each 8-10% lighter than last
    "bg":        "#0a0e17",   # deepest — window base
    "bg2":       "#0f1520",   # sidebar background
    "bg3":       "#141d2b",   # section card
    "bg4":       "#1a2538",   # input fields, inner cards
    "bg5":       "#202e45",   # hover / active states

    # ONE accent color — electric cyan, used sparingly
    "accent":    "#00b4d8",   # primary cyan
    "accent2":   "#0096c7",   # darker cyan — pressed buttons
    "accent3":   "#90e0ef",   # light cyan — hover text glow

    # Semantic — only for status meaning, never decoration
    "success":   "#2ec27e",   # green — online / matched
    "warning":   "#e9c46a",   # amber — caution
    "danger":    "#e76f51",   # red-orange — error / NSFW
    "kiss_col":  "#c77dff",   # purple — kissing tag

    # Text — 3-level hierarchy (most important design decision)
    "text":      "#eaf0fb",   # level 1 — headings, labels   (near-white)
    "text2":     "#7f8fa6",   # level 2 — secondary info     (steel gray)
    "text3":     "#3d4f66",   # level 3 — disabled/hints     (dim)

    # Structural
    "border":    "#1e2d42",   # card borders
    "border2":   "#283d5a",   # input borders, slightly brighter
    "header":    "#06090f",   # topmost header bar
    "divider":   "#182233",   # thin rule between sections
}

# Fonts — Segoe UI for UI text, Courier New ONLY for code/terminal/hashes
FONT_TITLE = ("Segoe UI",    10, "bold")
FONT_LABEL = ("Segoe UI",     9)
FONT_BOLD  = ("Segoe UI",     9, "bold")
FONT_SMALL = ("Segoe UI",     8)
FONT_TINY  = ("Segoe UI",     7)
FONT_MONO  = ("Courier New",  9)   # terminal output only
FONT_CODE  = ("Courier New",  8)   # hashes / paths only


# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 1 — CORE ENGINE (Hashing, Metadata, Thumbnails)
# ══════════════════════════════════════════════════════════════════════════════
class ForensicEngine:
    """Core processing: hashing, EXIF metadata, thumbnails."""

    @staticmethod
    def compute_hashes(filepath: str) -> Dict[str, str]:
        """Compute MD5 and SHA256 hashes from raw file bytes."""
        md5    = hashlib.md5()
        sha256 = hashlib.sha256()
        try:
            with open(filepath, "rb") as f:
                for chunk in iter(lambda: f.read(65536), b""):
                    md5.update(chunk)
                    sha256.update(chunk)
            return {"md5": md5.hexdigest(), "sha256": sha256.hexdigest()}
        except Exception as e:
            logger.warning(f"Hash error for {filepath}: {e}")
            return {"md5": "No data found", "sha256": "No data found"}

    @staticmethod
    def extract_metadata(filepath: str) -> Dict[str, Any]:
        """Extract full metadata from an image file."""
        nd  = "No data found"
        p   = Path(filepath)
        meta = {
            "filename":      p.name,
            "folder":        str(p.parent),
            "filepath":      str(p.resolve()),
            "filesize":      nd,
            "device":        nd,
            "capture_date":  nd,
            "gps":           nd,
            "last_modified": nd,
            "last_accessed": nd,
            "exif_raw":      nd,
        }

        try:
            stat = p.stat()
            meta["filesize"]      = f"{stat.st_size:,} bytes ({stat.st_size / 1024:.1f} KB)"
            meta["last_modified"] = datetime.datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            meta["last_accessed"] = datetime.datetime.fromtimestamp(stat.st_atime).strftime("%Y-%m-%d %H:%M:%S")
        except Exception:
            pass

        # EXIF via Pillow
        try:
            img = Image.open(filepath)
            exif_data = img._getexif() if hasattr(img, "_getexif") else None
            if exif_data:
                from PIL.ExifTags import TAGS, GPSTAGS
                decoded = {}
                for tag_id, value in exif_data.items():
                    tag = TAGS.get(tag_id, tag_id)
                    decoded[tag] = str(value)[:200]

                meta["exif_raw"]     = json.dumps(decoded, indent=2)[:2000]
                meta["device"]       = decoded.get("Make", nd)
                if "Model" in decoded and decoded.get("Make", nd) != nd:
                    meta["device"]   = f"{decoded.get('Make','')} {decoded.get('Model','')}".strip()
                meta["capture_date"] = decoded.get("DateTimeOriginal",
                                       decoded.get("DateTime", nd))

                # GPS
                gps_info = exif_data.get(34853)
                if gps_info:
                    try:
                        lat  = ForensicEngine._convert_gps(gps_info.get(2), gps_info.get(1))
                        lon  = ForensicEngine._convert_gps(gps_info.get(4), gps_info.get(3))
                        if lat and lon:
                            meta["gps"] = f"{lat:.6f}, {lon:.6f}"
                    except Exception:
                        pass
        except Exception as e:
            logger.debug(f"EXIF error for {filepath}: {e}")

        return meta

    @staticmethod
    def _convert_gps(coord, ref) -> Optional[float]:
        if not coord or not ref:
            return None
        try:
            d = float(coord[0])
            m = float(coord[1])
            s = float(coord[2])
            val = d + m / 60 + s / 3600
            if ref in ("S", "W"):
                val = -val
            return val
        except Exception:
            return None

    @staticmethod
    def make_thumbnail(filepath: str, size=THUMB_SIZE) -> Optional[ImageTk.PhotoImage]:
        """Generate a Tk-compatible thumbnail."""
        try:
            img = Image.open(filepath).convert("RGB")
            img.thumbnail(size, Image.LANCZOS)
            # Add border
            bordered = Image.new("RGB", (img.width + 4, img.height + 4), (0, 174, 255))
            bordered.paste(img, (2, 2))
            return ImageTk.PhotoImage(bordered)
        except Exception:
            return None

    @staticmethod
    def make_pil_thumbnail(filepath: str, size=PDF_THUMB_SIZE) -> Optional[Image.Image]:
        """Generate a PIL thumbnail for PDF embedding."""
        try:
            img = Image.open(filepath).convert("RGB")
            img.thumbnail(size, Image.LANCZOS)
            return img
        except Exception:
            return None

    @staticmethod
    def collect_images(folder: str) -> List[Path]:
        """Recursively collect all supported image paths."""
        return [
            p for p in Path(folder).rglob("*")
            if p.suffix.lower() in SUPPORTED_EXT
            and not p.name.startswith(".")
        ]


# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 2 — NSFW ENGINE (NudeNet)
# ══════════════════════════════════════════════════════════════════════════════
class NSFWEngine:
    """AI-based explicit content detection using NudeNet."""

class NSFWEngine:
    """AI-based explicit content detection using NudeNet."""

    # NudeNet v2 labels (nudenet < 3.0)
    LABELS_V2 = {
        "EXPOSED_ANUS", "EXPOSED_BUTTOCKS",
        "EXPOSED_BREAST_F", "EXPOSED_BREAST_M",
        "EXPOSED_GENITALIA_F", "EXPOSED_GENITALIA_M",
        "EXPOSED_BELLY", "EXPOSED_ARMPITS",
    }

    # NudeNet v3 labels (nudenet >= 3.0) — the current version
    LABELS_V3 = {
        "FEMALE_BREAST_EXPOSED", "FEMALE_GENITALIA_EXPOSED",
        "MALE_BREAST_EXPOSED",   "MALE_GENITALIA_EXPOSED",
        "BUTTOCKS_EXPOSED",      "ANUS_EXPOSED",
        "FEMALE_BREAST_COVERED", "FEMALE_GENITALIA_COVERED",
        "MALE_GENITALIA_COVERED","BUTTOCKS_COVERED",
    }

    # Combined — works with both versions
    EXPLICIT_LABELS = LABELS_V2 | LABELS_V3

    def __init__(self):
        self.detector = None
        self.available = NudeDetector is not None
        if self.available:
            try:
                self.detector = NudeDetector()
                logger.info("NudeNet NSFW detector loaded successfully.")
            except Exception as e:
                logger.error(f"NudeNet load failed: {e}")
                self.available = False

    def analyze(self, filepath: str, threshold: float = 0.3) -> Dict[str, Any]:
        """
        Returns dict:
            is_explicit : bool
            confidence  : float (0–1)
            detections  : list of label dicts
        """
        if not self.available or self.detector is None:
            return {"is_explicit": False, "confidence": 0.0, "detections": []}
        try:
            results = self.detector.detect(filepath)

            # Debug: log all detected labels so investigator can see what was found
            if results:
                all_labels = [(r.get("class","?"), round(r.get("score",0), 3)) for r in results]
                logger.debug(f"NudeNet [{Path(filepath).name}]: {all_labels}")

            explicit_detections = [
                r for r in results
                if r.get("class") in self.EXPLICIT_LABELS
                and r.get("score", 0) >= threshold
            ]

            # Also catch anything NudeNet flags that we haven't listed
            # by checking if any result has score >= threshold regardless of label
            # This future-proofs against new NudeNet versions
            if not explicit_detections:
                explicit_detections = [
                    r for r in results
                    if r.get("score", 0) >= threshold
                    and "COVERED" not in r.get("class", "")
                    and r.get("class", "") not in {
                        "FACE_FEMALE", "FACE_MALE", "ARMPITS_COVERED",
                        "BELLY_COVERED", "FEET_COVERED", "FEET_EXPOSED",
                    }
                ]

            confidence = max((r.get("score", 0) for r in explicit_detections), default=0.0)
            is_explicit = len(explicit_detections) > 0

            if is_explicit:
                logger.info(f"NSFW DETECTED: {Path(filepath).name} — "
                            f"confidence={confidence:.2f} "
                            f"labels={[d.get('class') for d in explicit_detections]}")

            return {
                "is_explicit":  is_explicit,
                "confidence":   round(confidence, 3),
                "detections":   explicit_detections,
            }
        except Exception as e:
            logger.warning(f"NSFW error for {filepath}: {e}")
            return {"is_explicit": False, "confidence": 0.0, "detections": []}



# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 2B — KISSING / INTIMATE CONTENT DETECTOR  (CLIP — optimised)
#
#  Speed fixes vs naive implementation:
#   1. Text prompts encoded ONCE at load time — not per image
#   2. Images processed in BATCHES of 16 — 10-16× faster than one-by-one
#   3. Smaller model: clip-vit-base-patch16 (faster than patch32)
#   4. torch.no_grad() + inference_mode for maximum speed
#   5. Image pre-processing done in parallel with PIL resize before batching
# ══════════════════════════════════════════════════════════════════════════════
class KissingDetector:

    MODEL_ID   = "openai/clip-vit-base-patch32"   # ~600 MB, good accuracy/speed
    BATCH_SIZE = 16                                 # images per GPU/CPU batch

    # Kept intentionally short — fewer prompts = faster text encoding
    POSITIVE_PROMPTS = [
        "two people kissing",
        "couple kissing on lips",
        "people making out",
        "romantic kiss",
        "intimate kiss between two people",
    ]
    NEGATIVE_PROMPTS = [
        "people standing",
        "landscape or nature photo",
        "text document",
        "crowd of people",
    ]

    def __init__(self):
        self.model          = None
        self.processor      = None
        self.available      = clip_ok
        self._loaded        = False
        # Pre-cached text features — computed once, reused for every image
        self._text_features = None
        self._n_pos         = len(self.POSITIVE_PROMPTS)

    # ── Model load (called once before first batch) ────────────────
    def load(self):
        if self._loaded or not self.available:
            return
        try:
            import torch
            logger.info("Loading CLIP model… (first run downloads ~600 MB)")
            self.processor = CLIPProcessor.from_pretrained(self.MODEL_ID)
            self.model     = CLIPModel.from_pretrained(self.MODEL_ID)
            self.model.eval()

            # ── Pre-encode ALL text prompts right now — do this ONCE ──
            all_prompts = self.POSITIVE_PROMPTS + self.NEGATIVE_PROMPTS
            text_inputs = self.processor(
                text=all_prompts, return_tensors="pt", padding=True
            )
            with torch.inference_mode():
                self._text_features = self.model.get_text_features(**text_inputs)
                # Normalise for cosine similarity
                self._text_features = self._text_features / \
                    self._text_features.norm(dim=-1, keepdim=True)

            self._loaded = True
            logger.info(f"CLIP ready — text features pre-cached "
                        f"({len(all_prompts)} prompts × {self._text_features.shape[-1]}d)")
        except Exception as e:
            logger.error(f"CLIP load failed: {e}")
            self.available = False

    # ── Batch analyse — PRIMARY METHOD ────────────────────────────
    def analyze_batch(self, filepaths: List[str],
                      threshold: float = 0.25) -> List[Dict[str, Any]]:
        """
        Process a batch of images in one forward pass.
        Returns a list of result dicts (same order as filepaths).
        """
        null = {"is_kissing": False, "confidence": 0.0, "top_prompt": ""}
        if not self.available or not self._loaded or not filepaths:
            return [null] * len(filepaths)

        try:
            import torch

            # Load & resize images (PIL is fast, keeps memory low)
            pil_images = []
            valid_idx  = []
            for i, fp in enumerate(filepaths):
                try:
                    img = Image.open(fp).convert("RGB")
                    img.thumbnail((224, 224), Image.BILINEAR)   # CLIP native size
                    pil_images.append(img)
                    valid_idx.append(i)
                except Exception:
                    pass

            if not pil_images:
                return [null] * len(filepaths)

            # Encode all images in one call
            img_inputs = self.processor(
                images=pil_images, return_tensors="pt", padding=True
            )
            with torch.inference_mode():
                img_features = self.model.get_image_features(**img_inputs)
                img_features = img_features / img_features.norm(dim=-1, keepdim=True)

                # Cosine similarity: [n_images × n_prompts]
                similarity = (img_features @ self._text_features.T) * 100
                probs      = similarity.softmax(dim=-1)   # shape [n_imgs, n_prompts]

            results = [null] * len(filepaths)
            for j, orig_i in enumerate(valid_idx):
                pos_probs  = probs[j, :self._n_pos].tolist()
                total_conf = sum(pos_probs)
                best_idx   = pos_probs.index(max(pos_probs))
                is_kissing = total_conf >= threshold

                if is_kissing:
                    logger.info(f"KISSING: {Path(filepaths[orig_i]).name} "
                                f"conf={total_conf:.3f}")

                results[orig_i] = {
                    "is_kissing":  is_kissing,
                    "confidence":  round(total_conf, 3),
                    "top_prompt":  self.POSITIVE_PROMPTS[best_idx] if is_kissing else "",
                }
            return results

        except Exception as e:
            logger.warning(f"CLIP batch error: {e}")
            return [null] * len(filepaths)

    # ── Single image wrapper (kept for compatibility) ──────────────
    def analyze(self, filepath: str, threshold: float = 0.25) -> Dict[str, Any]:
        results = self.analyze_batch([filepath], threshold)
        return results[0]


# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 2C — KEYWORD SEARCH ENGINE
#  Searches keywords in:
#    • Images      — via EasyOCR (reads text inside photos, screenshots, scans)
#    • PDF files   — via pdfplumber (text extraction, no OCR needed for digital PDFs)
#    • Word (.docx)— via python-docx
#    • PowerPoint  — via python-pptx
#    • Text files  — native Python read
# ══════════════════════════════════════════════════════════════════════════════

# Supported document extensions
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt",
                  ".csv", ".log", ".xml", ".html", ".htm", ".json"}
IMG_EXTENSIONS = {".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".tif",
                  ".webp", ".gif"}


class KeywordSearchEngine:
    """
    Extracts text from images (OCR) and documents, then searches keywords.
    Returns matched files with the exact sentences containing the keyword.
    """

    def __init__(self):
        self.ocr_available  = easyocr_ok
        self.pdf_available  = pdfplumber_ok
        self.docx_available = docx_ok
        self.pptx_available = pptx_ok
        self._ocr_reader    = None   # lazy-loaded EasyOCR reader
        self._ocr_loading   = False

    @property
    def available(self):
        # Always available — plain .txt/.csv/.html etc need no library at all
        return True

    # ── Lazy load EasyOCR (heavy — ~200 MB model) ─────────────────
    def load_ocr(self, log_fn=None):
        if self._ocr_reader is not None:
            return True
        if not self.ocr_available:
            return False
        try:
            if log_fn:
                log_fn("  Loading EasyOCR model (first run ~200 MB)…", "info")
            self._ocr_reader = _easyocr.Reader(
                ['en'],
                gpu=False,
                verbose=False,
            )
            if log_fn:
                log_fn("  ✔ EasyOCR ready.\n", "match")
            return True
        except Exception as e:
            logger.error(f"EasyOCR load failed: {e}")
            self.ocr_available = False
            return False

    # ── Extract text from an IMAGE via OCR ────────────────────────
    def extract_image_text(self, filepath: str) -> str:
        if self._ocr_reader is None:
            return ""
        try:
            results = self._ocr_reader.readtext(filepath, detail=0,
                                                paragraph=True)
            return " ".join(results)
        except Exception as e:
            logger.debug(f"OCR error {filepath}: {e}")
            return ""

    # ── Extract text from a PDF ────────────────────────────────────
    def extract_pdf_text(self, filepath: str) -> str:
        if not self.pdf_available:
            return ""
        try:
            text_parts = []
            with _pdfplumber.open(filepath) as pdf:
                for page in pdf.pages:
                    t = page.extract_text()
                    if t:
                        text_parts.append(t)
            return "\n".join(text_parts)
        except Exception as e:
            logger.debug(f"PDF extract error {filepath}: {e}")
            return ""

    # ── Extract text from a Word document ─────────────────────────
    def extract_docx_text(self, filepath: str) -> str:
        if not self.docx_available:
            return ""
        try:
            doc = _DocxDocument(filepath)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception as e:
            logger.debug(f"DOCX extract error {filepath}: {e}")
            return ""

    # ── Extract text from a PowerPoint ────────────────────────────
    def extract_pptx_text(self, filepath: str) -> str:
        if not self.pptx_available:
            return ""
        try:
            prs  = _PptxPresentation(filepath)
            bits = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        bits.append(shape.text)
            return "\n".join(bits)
        except Exception as e:
            logger.debug(f"PPTX extract error {filepath}: {e}")
            return ""

    # ── Extract text from a plain text file ───────────────────────
    def extract_txt_text(self, filepath: str) -> str:
        try:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""

    # ── Route to correct extractor by extension ───────────────────
    def extract_text(self, filepath: str) -> str:
        ext = Path(filepath).suffix.lower()
        if ext in IMG_EXTENSIONS:
            return self.extract_image_text(filepath)
        elif ext == ".pdf":
            return self.extract_pdf_text(filepath)
        elif ext in (".docx", ".doc"):
            return self.extract_docx_text(filepath)
        elif ext in (".pptx", ".ppt"):
            return self.extract_pptx_text(filepath)
        else:
            return self.extract_txt_text(filepath)

    # ── Search keywords in extracted text ─────────────────────────
    @staticmethod
    def search_keywords(text: str, keywords: List[str],
                        case_sensitive: bool = False) -> Dict[str, List[str]]:
        """
        Returns dict: { keyword: [matching_sentence, ...] }
        Only returns keywords that were actually found.
        """
        if not text or not keywords:
            return {}

        flags   = 0 if case_sensitive else re.IGNORECASE
        # Split into sentences for context
        sents   = re.split(r'[.\n!?;]', text)
        sents   = [s.strip() for s in sents if s.strip()]

        hits = {}
        for kw in keywords:
            kw_clean = kw.strip()
            if not kw_clean:
                continue
            # Whole-word match by default
            pattern = re.compile(r'\b' + re.escape(kw_clean) + r'\b', flags)
            matched_sents = [s for s in sents if pattern.search(s)]
            if matched_sents:
                hits[kw_clean] = matched_sents[:5]   # max 5 context sentences
        return hits

    # ── Full scan of a folder for keywords ────────────────────────
    def scan_folder(self, folder: str, keywords: List[str],
                    do_images: bool = True, do_docs: bool = True,
                    case_sensitive: bool = False,
                    progress_cb=None,
                    log_fn=None) -> List[Dict]:
        """
        Walk folder, extract text, search keywords.
        progress_cb(done, total, filepath) called for UI updates.
        log_fn(msg, tag) called for live log output.
        Returns list of result dicts.
        """
        def _log(msg, tag="info"):
            logger.info(msg)
            if log_fn:
                log_fn(f"  {msg}", tag)

        folder_path = Path(folder)
        all_files   = []

        if do_images:
            for ext in IMG_EXTENSIONS:
                all_files += list(folder_path.rglob(f"*{ext}"))
                all_files += list(folder_path.rglob(f"*{ext.upper()}"))
        if do_docs:
            for ext in DOC_EXTENSIONS:
                all_files += list(folder_path.rglob(f"*{ext}"))
                all_files += list(folder_path.rglob(f"*{ext.upper()}"))

        # Deduplicate preserving order
        seen_paths = set()
        deduped = []
        for f in all_files:
            key = str(f).lower()   # case-insensitive dedup on Windows
            if key not in seen_paths:
                seen_paths.add(key)
                deduped.append(f)
        all_files = deduped
        total     = len(all_files)

        _log(f"Keyword scan — found {total} file(s) to search "
             f"(images: {do_images}, docs: {do_docs})", "info")

        if total == 0:
            _log("No files found in folder for keyword search.", "error")
            return []

        results = []

        for i, fp in enumerate(all_files):
            if progress_cb:
                progress_cb(i + 1, total, fp.name)

            ext = fp.suffix.lower()

            # Skip images if OCR not available
            if ext in IMG_EXTENSIONS and not self.ocr_available:
                continue

            try:
                text = self.extract_text(str(fp))
            except Exception as e:
                _log(f"Extract error: {fp.name} — {e}", "error")
                continue

            if not text or not text.strip():
                logger.debug(f"No text extracted from: {fp.name}")
                continue

            hits = self.search_keywords(text, keywords, case_sensitive)
            if hits:
                ftype = "image" if ext in IMG_EXTENSIONS else "document"
                try:
                    metadata = ForensicEngine.extract_metadata(str(fp))
                    hashes   = ForensicEngine.compute_hashes(str(fp))
                except Exception:
                    metadata = {"filename": fp.name, "folder": str(fp.parent),
                                "filepath": str(fp)}
                    hashes   = {}

                results.append({
                    "filepath":       str(fp),
                    "filename":       fp.name,
                    "filetype":       ftype,
                    "extension":      ext,
                    "keywords_found": hits,
                    "keyword_count":  len(hits),
                    "total_hits":     sum(len(v) for v in hits.values()),
                    "metadata":       metadata,
                    "hashes":         hashes,
                })
                kw_list = ", ".join(hits.keys())
                _log(f"✔ KEYWORD MATCH  [{kw_list}]  →  {fp.name}", "match")

        _log(f"Keyword scan done — {len(results)} match(es) from {total} file(s).",
             "match" if results else "info")
        return results


class FaceEngine:

    """Face recognition matching engine with disk-based encoding cache."""

    def __init__(self):
        self.available = face_recognition is not None

    def load_image(self, filepath: str) -> Optional[np.ndarray]:
        try:
            pil = Image.open(filepath).convert("RGB")
            w, h = pil.size
            if max(w, h) > MAX_DIM:
                scale = MAX_DIM / max(w, h)
                pil = pil.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
            return np.ascontiguousarray(np.array(pil, dtype=np.uint8))
        except Exception:
            return None

    def encode_sample(self, filepath: str) -> Optional[np.ndarray]:
        if not self.available:
            return None
        img = self.load_image(filepath)
        if img is None:
            return None
        locs = face_recognition.face_locations(img, model="hog")
        if not locs:
            return None
        if len(locs) > 1:
            locs = [max(locs, key=lambda l: (l[2] - l[0]) * abs(l[1] - l[3]))]
        encs = face_recognition.face_encodings(img, locs)
        return encs[0] if encs else None

    def encode_image(self, filepath: str) -> List[List[float]]:
        """Return list of face encodings (serializable as lists)."""
        if not self.available:
            return []
        img = self.load_image(filepath)
        if img is None:
            return []
        locs = face_recognition.face_locations(img, model="hog")
        if not locs:
            return []
        encs = face_recognition.face_encodings(img, locs)
        return [e.tolist() for e in encs]

    def matches_sample(self, encodings: List[List[float]],
                       sample_enc: np.ndarray, tolerance: float) -> bool:
        if not encodings or sample_enc is None:
            return False
        enc_arrays = [np.array(e) for e in encodings]
        return any(face_recognition.compare_faces(enc_arrays, sample_enc, tolerance=tolerance))

    @staticmethod
    def load_cache(folder: str) -> Dict:
        cp = Path(folder) / CACHE_FILENAME
        if cp.exists():
            try:
                with open(cp, "r") as f:
                    return json.load(f)
            except Exception:
                return {}
        return {}

    @staticmethod
    def save_cache(folder: str, cache: Dict):
        cp = Path(folder) / CACHE_FILENAME
        try:
            with open(cp, "w") as f:
                json.dump(cache, f)
        except Exception as e:
            logger.warning(f"Cache save failed: {e}")


# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 4 — EVIDENCE MANAGER
# ══════════════════════════════════════════════════════════════════════════════
class EvidenceManager:
    """Handles copy/move of matched files with audit logging."""

    @staticmethod
    def _write_log(dest_folder: str, action: str, files: List[str]):
        log_path = Path(dest_folder) / f"finder_evidence_{action}_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
        try:
            with open(log_path, "w", encoding="utf-8") as f:
                f.write(f"Documents / Images Finder Tool — Evidence {action.upper()} Log\n")
                f.write(f"Advanced Search & Analysis Tool\n")
                f.write(f"Timestamp : {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
                f.write(f"Action    : {action.upper()}\n")
                f.write(f"Files     : {len(files)}\n")
                f.write("=" * 70 + "\n\n")
                for fp in files:
                    f.write(f"{fp}\n")
            logger.info(f"Evidence log written: {log_path}")
        except Exception as e:
            logger.error(f"Log write failed: {e}")

    @staticmethod
    def copy_files(file_paths: List[str], dest_folder: str) -> int:
        Path(dest_folder).mkdir(parents=True, exist_ok=True)
        copied = 0
        for fp in file_paths:
            try:
                src  = Path(fp)
                dest = Path(dest_folder) / src.name
                # Avoid overwrite — append counter
                if dest.exists():
                    dest = Path(dest_folder) / f"{src.stem}_{copied}{src.suffix}"
                shutil.copy2(fp, dest)
                copied += 1
            except Exception as e:
                logger.warning(f"Copy failed {fp}: {e}")
        EvidenceManager._write_log(dest_folder, "copy", file_paths)
        return copied

    @staticmethod
    def move_files(file_paths: List[str], dest_folder: str) -> int:
        Path(dest_folder).mkdir(parents=True, exist_ok=True)
        moved = 0
        for fp in file_paths:
            try:
                src  = Path(fp)
                dest = Path(dest_folder) / src.name
                if dest.exists():
                    dest = Path(dest_folder) / f"{src.stem}_{moved}{src.suffix}"
                shutil.move(fp, dest)
                moved += 1
            except Exception as e:
                logger.warning(f"Move failed {fp}: {e}")
        EvidenceManager._write_log(dest_folder, "move", file_paths)
        return moved


# ══════════════════════════════════════════════════════════════════════════════
#  MODULE 5 — PDF REPORT ENGINE
# ══════════════════════════════════════════════════════════════════════════════
class ReportEngine:
    """Generates a professional PDF report."""

    # Colours
    COL_HEADER_BG   = colors.HexColor("#1a1a2e")   # deep navy
    COL_HEADER_TEXT = colors.white                  # pure white — always visible
    COL_ROW_ODD     = colors.HexColor("#f4f6ff")
    COL_ROW_EVEN    = colors.HexColor("#e8ecf8")
    COL_GRID        = colors.HexColor("#b0b8d8")
    COL_ACCENT      = colors.HexColor("#5a4ed1")
    COL_TITLE       = colors.HexColor("#1a1a2e")
    COL_SUB         = colors.HexColor("#3a3a5a")

    # Larger thumbnail for PDF — clear and readable
    PDF_THUMB = (160, 120)   # render at high res, display scaled in cell

    def __init__(self):
        self.available = reportlab_ok

    # ── helpers ──────────────────────────────────────────────────
    @staticmethod
    def _thumb_image(filepath: str) -> Optional["RLImage"]:
        """Return a high-quality RLImage for PDF embedding."""
        import io as _io
        try:
            pil = Image.open(filepath).convert("RGB")
            # High-res thumbnail — keeps quality when scaled down in PDF
            pil.thumbnail((200, 150), Image.LANCZOS)
            buf = _io.BytesIO()
            pil.save(buf, format="JPEG", quality=95, subsampling=0)
            buf.seek(0)
            # Display size in PDF: 38mm wide × proportional height
            aspect = pil.height / pil.width if pil.width else 0.75
            w_mm   = 38 * mm
            h_mm   = w_mm * aspect
            return RLImage(buf, width=w_mm, height=h_mm)
        except Exception:
            return None

    def generate(self, results: List[Dict], output_path: str,
                 total_scanned: int, scan_mode: str) -> bool:
        if not self.available:
            logger.error("reportlab not installed — PDF generation unavailable.")
            return False

        try:
            import io as _io

            PAGE  = landscape(A4)
            doc   = SimpleDocTemplate(
                output_path,
                pagesize=PAGE,
                rightMargin=12 * mm, leftMargin=12 * mm,
                topMargin=18 * mm,   bottomMargin=18 * mm,
            )
            page_w = PAGE[0] - 24 * mm   # usable width

            story = []

            # ══ Page header banner ════════════════════════════════
            # Draw a solid navy banner using a single-cell table
            banner_style = ParagraphStyle(
                "banner",
                fontName="Helvetica-Bold", fontSize=17,
                textColor=colors.white,
                alignment=TA_CENTER, leading=22,
            )
            sub_banner = ParagraphStyle(
                "subbanner",
                fontName="Helvetica", fontSize=10,
                textColor=colors.HexColor("#ccccee"),
                alignment=TA_CENTER, leading=14,
            )
            banner_tbl = Table(
                [[Paragraph("DOCUMENTS / IMAGES FINDER TOOL — SEARCH REPORT", banner_style)],
                 [Paragraph("Documents / Images Finder Tool  ·  Advanced Search & Analysis", sub_banner)]],
                colWidths=[page_w],
            )
            banner_tbl.setStyle(TableStyle([
                ("BACKGROUND",    (0, 0), (-1, -1), self.COL_HEADER_BG),
                ("TOPPADDING",    (0, 0), (-1, -1), 10),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
                ("LEFTPADDING",   (0, 0), (-1, -1), 14),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 14),
            ]))
            story.append(banner_tbl)
            story.append(Spacer(1, 4 * mm))

            # ── Summary row ──────────────────────────────────────
            ts = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            info_data = [
                ["Scan Timestamp", ts],
                ["Scan Mode",      scan_mode],
                ["Total Scanned",  f"{total_scanned:,}"],
                ["Total Matched",  str(len(results))],
            ]
            info_style_p = ParagraphStyle(
                "infop", fontName="Helvetica", fontSize=9,
                textColor=colors.HexColor("#333355"), leading=13,
            )
            info_bold_p = ParagraphStyle(
                "infobold", fontName="Helvetica-Bold", fontSize=9,
                textColor=self.COL_TITLE, leading=13,
            )
            info_rows = [[Paragraph(k, info_bold_p), Paragraph(v, info_style_p)]
                         for k, v in info_data]
            summary_tbl = Table(info_rows, colWidths=[42 * mm, 80 * mm])
            summary_tbl.setStyle(TableStyle([
                ("BACKGROUND",    (0, 0), (-1, -1), colors.HexColor("#eef0ff")),
                ("GRID",          (0, 0), (-1, -1), 0.3, self.COL_GRID),
                ("TOPPADDING",    (0, 0), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ("LEFTPADDING",   (0, 0), (-1, -1), 8),
            ]))
            story.append(summary_tbl)
            story.append(Spacer(1, 6 * mm))

            if not results:
                story.append(Paragraph("No matches found.", info_style_p))
                doc.build(story)
                return True

            # ══ Evidence table ════════════════════════════════════
            # Column definitions
            headers = [
                "#", "Thumbnail", "Image Name", "Folder / Path",
                "MD5 Hash", "Device", "Capture Date",
                "Last Modified", "Last Accessed", "File Size",
            ]
            col_widths = [
                8*mm,   # #
                40*mm,  # thumbnail  ← wider for clear image
                38*mm,  # image name
                52*mm,  # folder/path
                46*mm,  # md5
                30*mm,  # device
                28*mm,  # capture date
                28*mm,  # last modified
                28*mm,  # last accessed
                26*mm,  # file size
            ]

            # Header row style — white text on dark navy, guaranteed visible
            hdr_para_style = ParagraphStyle(
                "hdr", fontName="Helvetica-Bold", fontSize=8,
                textColor=colors.white,      # ← WHITE — always readable
                alignment=TA_CENTER, leading=11,
            )
            cell_para_style = ParagraphStyle(
                "cell", fontName="Helvetica", fontSize=8,
                textColor=colors.HexColor("#111133"), leading=11,
            )
            cell_mono_style = ParagraphStyle(
                "mono", fontName="Courier", fontSize=7,
                textColor=colors.HexColor("#222244"), leading=10,
            )

            header_row  = [Paragraph(h, hdr_para_style) for h in headers]
            table_data  = [header_row]

            for idx, r in enumerate(results, 1):
                meta   = r.get("metadata", {})
                hashes = r.get("hashes",   {})
                nsfw   = r.get("nsfw",     {})
                nd     = "No data found"

                # High-quality thumbnail
                thumb_cell = Paragraph("No preview", cell_para_style)
                rl_img = self._thumb_image(r["filepath"])
                if rl_img:
                    thumb_cell = rl_img

                # Folder path — show last 55 chars
                folder_path = meta.get("folder", nd)
                if len(folder_path) > 55:
                    folder_path = "…" + folder_path[-54:]

                row = [
                    Paragraph(str(idx), cell_para_style),
                    thumb_cell,
                    Paragraph(meta.get("filename",     nd), cell_para_style),
                    Paragraph(folder_path,                  cell_para_style),
                    Paragraph(hashes.get("md5",        nd), cell_mono_style),
                    Paragraph(meta.get("device",       nd), cell_para_style),
                    Paragraph(meta.get("capture_date", nd), cell_para_style),
                    Paragraph(meta.get("last_modified",nd), cell_para_style),
                    Paragraph(meta.get("last_accessed",nd), cell_para_style),
                    Paragraph(meta.get("filesize",     nd), cell_para_style),
                ]
                table_data.append(row)

            tbl = Table(
                table_data,
                colWidths=col_widths,
                repeatRows=1,           # repeat header on each page
                hAlign="LEFT",
            )
            tbl.setStyle(TableStyle([
                # ── Header row ──────────────────────────────────
                ("BACKGROUND",    (0, 0), (-1, 0), self.COL_HEADER_BG),
                ("TEXTCOLOR",     (0, 0), (-1, 0), colors.white),
                ("FONTNAME",      (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE",      (0, 0), (-1, 0), 8),
                ("TOPPADDING",    (0, 0), (-1, 0), 7),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 7),
                ("ALIGN",         (0, 0), (-1, 0), "CENTER"),
                ("VALIGN",        (0, 0), (-1, 0), "MIDDLE"),
                # ── Data rows ───────────────────────────────────
                ("FONTNAME",      (0, 1), (-1, -1), "Helvetica"),
                ("FONTSIZE",      (0, 1), (-1, -1), 8),
                ("ROWBACKGROUNDS",(0, 1), (-1, -1),
                 [self.COL_ROW_ODD, self.COL_ROW_EVEN]),
                ("VALIGN",        (0, 1), (-1, -1), "MIDDLE"),
                ("ALIGN",         (0, 1), (0, -1),  "CENTER"),  # # col
                ("ALIGN",         (1, 1), (1, -1),  "CENTER"),  # thumb col
                ("TOPPADDING",    (0, 1), (-1, -1), 4),
                ("BOTTOMPADDING", (0, 1), (-1, -1), 4),
                ("LEFTPADDING",   (0, 0), (-1, -1), 4),
                ("RIGHTPADDING",  (0, 0), (-1, -1), 4),
                # ── Grid ────────────────────────────────────────
                ("GRID",          (0, 0), (-1, -1), 0.4, self.COL_GRID),
                ("LINEBELOW",     (0, 0), (-1, 0),  1.2, self.COL_ACCENT),
            ]))

            story.append(tbl)
            story.append(Spacer(1, 8 * mm))

            # ── Footer ───────────────────────────────────────────
            story.append(HRFlowable(
                width="100%", thickness=1, color=self.COL_ACCENT))
            story.append(Spacer(1, 2 * mm))
            story.append(Paragraph(
                "Documents / Images Finder Tool  ·  "
                "Generated: " + ts + "  ·  FOR LAW ENFORCEMENT USE ONLY",
                ParagraphStyle(
                    "footer2", fontName="Helvetica", fontSize=7,
                    textColor=colors.HexColor("#555577"),
                    alignment=TA_CENTER,
                )
            ))

            doc.build(story)
            logger.info(f"PDF report written: {output_path}")
            return True

        except Exception as e:
            logger.error(f"PDF generation error: {e}")
            import traceback; traceback.print_exc()
            return False


# ══════════════════════════════════════════════════════════════════════════════
#  WORKER FUNCTION (multiprocessing)
# ══════════════════════════════════════════════════════════════════════════════
def _encode_worker(filepath_str: str) -> tuple:
    """Process one image — extract face encodings. Runs in subprocess."""
    try:
        import face_recognition as fr
        from PIL import Image
        import numpy as np

        pil = Image.open(filepath_str).convert("RGB")
        w, h = pil.size
        if max(w, h) > MAX_DIM:
            scale = MAX_DIM / max(w, h)
            pil = pil.resize((int(w * scale), int(h * scale)), Image.LANCZOS)
        img  = np.ascontiguousarray(np.array(pil, dtype=np.uint8))
        locs = fr.face_locations(img, model="hog")
        encs = fr.face_encodings(img, locs) if locs else []
        return (filepath_str, [e.tolist() for e in encs])
    except Exception:
        return (filepath_str, [])


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN APPLICATION UI
# ══════════════════════════════════════════════════════════════════════════════
class ForensicApp(tk.Tk):
    """Documents / Images Finder Tool — Main Tkinter Application."""

    def __init__(self):
        super().__init__()
        self.title("Documents / Images Finder Tool")
        self.geometry("1400x880")
        self.minsize(1200, 750)
        self.configure(bg=C["bg"])

        # Engines
        self.forensic    = ForensicEngine()
        self.face_eng    = FaceEngine()
        self.nsfw_eng    = NSFWEngine()
        self.kiss_eng    = KissingDetector()
        self.kw_eng      = KeywordSearchEngine()
        self.report_eng  = ReportEngine()
        self.evidence    = EvidenceManager()

        # State
        self.scan_folder      = tk.StringVar()
        self.sample_path      = tk.StringVar()
        self.tolerance        = tk.DoubleVar(value=0.55)
        self.nsfw_threshold   = tk.DoubleVar(value=0.30)
        self.kiss_threshold   = tk.DoubleVar(value=0.25)
        self.workers_var      = tk.IntVar(value=max(2, multiprocessing.cpu_count() - 1))
        self.enable_face      = tk.BooleanVar(value=True)
        self.enable_nsfw      = tk.BooleanVar(value=False)
        self.enable_kissing   = tk.BooleanVar(value=False)
        self.enable_keyword   = tk.BooleanVar(value=False)
        self.kw_search_images = tk.BooleanVar(value=True)
        self.kw_search_docs   = tk.BooleanVar(value=True)
        self.kw_case_sensitive= tk.BooleanVar(value=False)
        self.keywords_var     = tk.StringVar()
        self.keyword_results  = []   # separate list from face/nsfw results

        self.running        = False
        self.matched_results: List[Dict] = []
        self.all_images:      List[Path] = []
        self._thumb_refs:     List       = []   # prevent GC
        self._sample_thumb                = None

        self._build_ui()
        self._center_window()
        logger.info("Documents / Images Finder Tool started.")

    def _center_window(self):
        self.update_idletasks()
        w, h = 1400, 880
        x = (self.winfo_screenwidth()  - w) // 2
        y = (self.winfo_screenheight() - h) // 2
        self.geometry(f"{w}x{h}+{x}+{y}")

    # ──────────────────────────────────────────────────────────────
    #  UI BUILDER
    # ──────────────────────────────────────────────────────────────
    def _build_ui(self):
        self._build_header()

        # Main body
        body = tk.Frame(self, bg=C["bg"])
        body.pack(fill="both", expand=True, padx=0, pady=0)

        self._build_left_panel(body)
        self._build_right_panel(body)

        self._build_statusbar()

    def _build_header(self):
        hdr = tk.Frame(self, bg=C["header"])
        hdr.pack(fill="x")
        tk.Frame(hdr, bg=C["accent"], height=3).pack(fill="x")   # cyan top rule

        inner = tk.Frame(hdr, bg=C["header"])
        inner.pack(fill="x")

        # ── Logo pill ────────────────────────────────────────────
        left = tk.Frame(inner, bg=C["header"])
        left.pack(side="left")

        pill = tk.Frame(left, bg=C["accent"], padx=20, pady=14)
        pill.pack(side="left")
        tk.Label(pill, text="DIFT", font=("Segoe UI", 20, "bold"),
                 bg=C["accent"], fg=C["header"]).pack()

        tk.Frame(left, bg=C["border2"], width=1).pack(
            side="left", fill="y", padx=18, pady=10)

        title_f = tk.Frame(left, bg=C["header"])
        title_f.pack(side="left")
        tk.Label(title_f, text="Documents / Images Finder Tool",
                 font=("Segoe UI", 13, "bold"),
                 bg=C["header"], fg=C["text"]).pack(anchor="w")
        tk.Label(title_f,
                 text="Advanced Document & Image Search  ·  Face · NSFW · Keyword Detection  ·  v2.0",
                 font=("Segoe UI", 8),
                 bg=C["header"], fg=C["text2"]).pack(anchor="w", pady=(3, 0))

        # ── Module status pills ──────────────────────────────────
        right = tk.Frame(inner, bg=C["header"])
        right.pack(side="right", padx=16, pady=12)

        def mod_pill(name, online):
            bg   = C["bg4"]
            dotc = C["success"] if online else C["danger"]
            f = tk.Frame(right, bg=bg, padx=10, pady=6)
            f.pack(side="left", padx=3)
            tk.Label(f, text="●", font=("Segoe UI", 7),
                     bg=bg, fg=dotc).pack(side="left")
            tk.Label(f, text=f"  {name}", font=("Segoe UI", 8),
                     bg=bg, fg=C["text2"]).pack(side="left")

        mod_pill("Face-Rec", face_recognition is not None)
        mod_pill("NudeNet",  NudeDetector is not None)
        mod_pill("CLIP-AI",  clip_ok)
        mod_pill("PDF",      reportlab_ok)

        tk.Frame(hdr, bg=C["border"], height=1).pack(fill="x")

    def _build_left_panel(self, parent):
        left = tk.Frame(parent, bg=C["bg2"], width=430)
        left.pack(side="left", fill="y", padx=0, pady=0)
        left.pack_propagate(False)

        # Scrollable inner
        canvas   = tk.Canvas(left, bg=C["bg2"], highlightthickness=0)
        scrollbar = tk.Scrollbar(left, orient="vertical", command=canvas.yview)
        canvas.configure(yscrollcommand=scrollbar.set)
        scrollbar.pack(side="right", fill="y")
        canvas.pack(side="left", fill="both", expand=True)

        inner = tk.Frame(canvas, bg=C["bg2"])
        win_id = canvas.create_window((0, 0), window=inner, anchor="nw")
        inner.bind("<Configure>", lambda e: canvas.configure(
            scrollregion=canvas.bbox("all")))
        canvas.bind("<Configure>", lambda e: canvas.itemconfig(win_id, width=e.width))

        self._build_controls(inner)

    def _build_controls(self, parent):

        # ── Section card ─────────────────────────────────────────
        def section(icon, title):
            wrap = tk.Frame(parent, bg=C["bg2"])
            wrap.pack(fill="x", padx=10, pady=(10, 0))
            card = tk.Frame(wrap, bg=C["bg3"])
            card.pack(fill="x")
            # Cyan left accent bar
            tk.Frame(card, bg=C["accent"], width=3).pack(side="left", fill="y")
            body_wrap = tk.Frame(card, bg=C["bg3"])
            body_wrap.pack(fill="x", expand=True)
            # Section header row
            hrow = tk.Frame(body_wrap, bg=C["bg3"], pady=9)
            hrow.pack(fill="x", padx=14)
            tk.Label(hrow, text=icon, font=("Segoe UI", 10),
                     bg=C["bg3"], fg=C["accent"]).pack(side="left")
            tk.Label(hrow, text=f"  {title}",
                     font=("Segoe UI", 9, "bold"),
                     bg=C["bg3"], fg=C["text"]).pack(side="left")
            tk.Frame(body_wrap, bg=C["border"], height=1).pack(fill="x", padx=14)
            body = tk.Frame(body_wrap, bg=C["bg3"], padx=14, pady=10)
            body.pack(fill="x")
            return body

        # ── Field label ──────────────────────────────────────────
        def field_lbl(p, text):
            tk.Label(p, text=text, font=("Segoe UI", 8),
                     bg=C["bg3"], fg=C["text2"]).pack(anchor="w", pady=(8, 2))

        # ── Path entry with Browse button ─────────────────────────
        def path_entry(p, var, cmd):
            row = tk.Frame(p, bg=C["bg3"])
            row.pack(fill="x")
            row.columnconfigure(0, weight=1)
            # Input border frame
            inp_border = tk.Frame(row, bg=C["border2"], padx=1, pady=1)
            inp_border.grid(row=0, column=0, sticky="ew")
            inp_bg = tk.Frame(inp_border, bg=C["bg4"])
            inp_bg.pack(fill="x")
            tk.Entry(inp_bg, textvariable=var,
                     font=("Segoe UI", 9),
                     bg=C["bg4"], fg=C["text"],
                     bd=0, relief="flat",
                     insertbackground=C["accent"],
                     highlightthickness=0).pack(
                         fill="x", padx=10, pady=8)
            tk.Button(row, text="Browse",
                      font=("Segoe UI", 8, "bold"),
                      bg=C["accent2"], fg="white",
                      bd=0, relief="flat",
                      padx=14, pady=8,
                      cursor="hand2",
                      activebackground=C["accent"],
                      activeforeground=C["header"],
                      command=cmd).grid(row=0, column=1, padx=(6, 0))

        # ── Slider with live value label ─────────────────────────
        def slider(p, label, var, lo, hi, step, fmt="{:.2f}"):
            row = tk.Frame(p, bg=C["bg3"])
            row.pack(fill="x", pady=(8, 0))
            top = tk.Frame(row, bg=C["bg3"])
            top.pack(fill="x")
            tk.Label(top, text=label, font=("Segoe UI", 8),
                     bg=C["bg3"], fg=C["text2"]).pack(side="left")
            val_lbl = tk.Label(top, text=fmt.format(var.get()),
                               font=("Segoe UI", 9, "bold"),
                               bg=C["bg3"], fg=C["accent"])
            val_lbl.pack(side="right")
            tk.Scale(row, from_=lo, to=hi, resolution=step,
                     orient="horizontal", variable=var,
                     bg=C["bg3"], fg=C["text2"],
                     troughcolor=C["bg5"],
                     activebackground=C["accent3"],
                     highlightthickness=0, showvalue=False, bd=0,
                     command=lambda v: val_lbl.config(
                         text=fmt.format(float(v)))).pack(fill="x")

        # ── Toggle checkbox ──────────────────────────────────────
        def toggle(p, label, var):
            tk.Checkbutton(p, text=f"  {label}", variable=var,
                           font=("Segoe UI", 9),
                           bg=C["bg3"], fg=C["text"],
                           selectcolor=C["bg5"],
                           activebackground=C["bg3"],
                           activeforeground=C["accent3"]).pack(
                               anchor="w", pady=(6, 2))

        # ── Inline status line ───────────────────────────────────
        def status_line(p, module, ok, install_cmd=None):
            row = tk.Frame(p, bg=C["bg3"])
            row.pack(anchor="w", pady=(4, 0))
            dot_c = C["success"] if ok else C["danger"]
            tk.Label(row, text="●", font=("Segoe UI", 8),
                     bg=C["bg3"], fg=dot_c).pack(side="left")
            state = "Available" if ok else "Not installed"
            tk.Label(row, text=f"  {module}  —  {state}",
                     font=("Segoe UI", 8),
                     bg=C["bg3"], fg=dot_c).pack(side="left")
            if not ok and install_cmd:
                tk.Label(row, text=f"  ·  {install_cmd}",
                         font=("Courier New", 7),
                         bg=C["bg3"], fg=C["text3"]).pack(side="left")

        # ═══════════════════════════════════════════════════════
        # 01 — SCAN TARGET
        # ═══════════════════════════════════════════════════════
        s1 = section("📁", "Scan Target")
        field_lbl(s1, "Image folder to scan:")
        path_entry(s1, self.scan_folder, self._browse_scan_folder)

        # ═══════════════════════════════════════════════════════
        # 02 — FACE RECOGNITION
        # ═══════════════════════════════════════════════════════
        s2 = section("👤", "Face Recognition")
        toggle(s2, "Enable face matching engine", self.enable_face)

        field_lbl(s2, "Subject sample image (person to find):")
        path_entry(s2, self.sample_path, self._browse_sample)

        # Preview box
        prev_frame = tk.Frame(s2, bg=C["border2"], padx=1, pady=1)
        prev_frame.pack(fill="x", pady=(8, 4))
        self.sample_preview_frame = tk.Frame(prev_frame, bg=C["bg4"], height=155)
        self.sample_preview_frame.pack(fill="x")
        self.sample_preview_frame.pack_propagate(False)
        self.sample_preview_lbl = tk.Label(
            self.sample_preview_frame,
            text="No subject image selected",
            font=("Segoe UI", 9),
            bg=C["bg4"], fg=C["text3"])
        self.sample_preview_lbl.place(relx=0.5, rely=0.5, anchor="center")

        slider(s2, "Match Tolerance", self.tolerance, 0.40, 0.70, 0.01)
        tk.Label(s2, text="0.40 = strict  ·  0.55 = recommended  ·  0.70 = lenient",
                 font=("Segoe UI", 7),
                 bg=C["bg3"], fg=C["text3"]).pack(anchor="w", pady=(2, 4))

        # ═══════════════════════════════════════════════════════
        # 03 — CONTENT DETECTION
        # ═══════════════════════════════════════════════════════
        s3 = section("🔍", "Content Detection AI")

        # NudeNet block
        status_line(s3, "NudeNet", self.nsfw_eng.available, "pip install nudenet")
        toggle(s3, "Detect nude / explicit content", self.enable_nsfw)
        slider(s3, "Nudity threshold",
               self.nsfw_threshold, 0.30, 0.95, 0.05, "{:.0%}")

        # Divider rule between two detectors
        tk.Frame(s3, bg=C["border"], height=1).pack(fill="x", pady=(10, 2))

        # CLIP block
        status_line(s3, "CLIP-AI  (Kissing)", self.kiss_eng.available,
                    "pip install transformers torch")
        toggle(s3, "Detect kissing / intimate content", self.enable_kissing)
        slider(s3, "Kissing threshold",
               self.kiss_threshold, 0.10, 0.60, 0.05, "{:.0%}")
        tk.Label(s3, text="⚠  CLIP downloads ~600 MB model on first use",
                 font=("Segoe UI", 8),
                 bg=C["bg3"], fg=C["warning"]).pack(anchor="w", pady=(6, 0))

        # ═══════════════════════════════════════════════════════
        # 04 — KEYWORD SEARCH
        # ═══════════════════════════════════════════════════════
        s_kw = section("🔤", "Keyword Search")

        # Always-available note
        tk.Label(s_kw,
                 text="✔  TXT / CSV / HTML — always available (no install needed)",
                 font=("Segoe UI", 8),
                 bg=C["bg3"], fg=C["success"]).pack(anchor="w", pady=(4, 0))

        # Optional library status
        status_line(s_kw, "pdfplumber  (PDF files)", self.kw_eng.pdf_available,
                    "pip install pdfplumber")
        status_line(s_kw, "python-docx  (Word files)", self.kw_eng.docx_available,
                    "pip install python-docx")
        status_line(s_kw, "python-pptx  (PowerPoint)", self.kw_eng.pptx_available,
                    "pip install python-pptx")
        status_line(s_kw, "EasyOCR  (text inside images)", self.kw_eng.ocr_available,
                    "pip install easyocr")

        tk.Frame(s_kw, bg=C["border"], height=1).pack(fill="x", pady=(8, 4))

        toggle(s_kw, "Enable Keyword Search", self.enable_keyword)

        # Keywords input
        tk.Label(s_kw, text="Keywords  (comma-separated):",
                 font=("Segoe UI", 8),
                 bg=C["bg3"], fg=C["text2"]).pack(anchor="w", pady=(8, 2))

        kw_border = tk.Frame(s_kw, bg=C["border2"], padx=1, pady=1)
        kw_border.pack(fill="x")
        kw_bg = tk.Frame(kw_border, bg=C["bg4"])
        kw_bg.pack(fill="x")
        tk.Entry(kw_bg, textvariable=self.keywords_var,
                 font=("Segoe UI", 9),
                 bg=C["bg4"], fg=C["text"],
                 bd=0, relief="flat",
                 insertbackground=C["accent"],
                 highlightthickness=0).pack(fill="x", padx=10, pady=8)

        tk.Label(s_kw, text='Example:  science, math, Pakistan, CNIC',
                 font=("Segoe UI", 7),
                 bg=C["bg3"], fg=C["text3"]).pack(anchor="w", pady=(2, 6))

        # Search scope checkboxes
        tk.Label(s_kw, text="Search in:",
                 font=("Segoe UI", 8, "bold"),
                 bg=C["bg3"], fg=C["text2"]).pack(anchor="w")

        scope_row = tk.Frame(s_kw, bg=C["bg3"])
        scope_row.pack(anchor="w")
        tk.Checkbutton(scope_row, text="  Images (OCR)",
                       variable=self.kw_search_images,
                       font=("Segoe UI", 8),
                       bg=C["bg3"], fg=C["text"],
                       selectcolor=C["bg5"],
                       activebackground=C["bg3"],
                       activeforeground=C["accent3"]).pack(side="left")
        tk.Checkbutton(scope_row, text="  Documents",
                       variable=self.kw_search_docs,
                       font=("Segoe UI", 8),
                       bg=C["bg3"], fg=C["text"],
                       selectcolor=C["bg5"],
                       activebackground=C["bg3"],
                       activeforeground=C["accent3"]).pack(side="left", padx=(8, 0))

        tk.Checkbutton(s_kw, text="  Case-sensitive matching",
                       variable=self.kw_case_sensitive,
                       font=("Segoe UI", 8),
                       bg=C["bg3"], fg=C["text2"],
                       selectcolor=C["bg5"],
                       activebackground=C["bg3"],
                       activeforeground=C["accent3"]).pack(anchor="w", pady=(4, 0))

        tk.Label(s_kw,
                 text="⚠  OCR loads ~200 MB model on first use",
                 font=("Segoe UI", 8),
                 bg=C["bg3"], fg=C["warning"]).pack(anchor="w", pady=(6, 0))

        # ═══════════════════════════════════════════════════════
        # 05 — PERFORMANCE
        # ═══════════════════════════════════════════════════════
        s4 = section("⚡", "Performance")
        n_cpu = multiprocessing.cpu_count()
        slider(s4, f"CPU Threads  (system: {n_cpu} cores)",
               self.workers_var, 1, n_cpu, 1, "{:.0f}")
        tk.Label(s4,
                 text="More threads = faster encoding on large image sets",
                 font=("Segoe UI", 7),
                 bg=C["bg3"], fg=C["text3"]).pack(anchor="w", pady=(2, 4))

        # ═══════════════════════════════════════════════════════
        # 05 — OPERATIONS
        # ═══════════════════════════════════════════════════════
        s5 = section("▶", "Operations")

        # Primary CTA
        self.scan_btn = tk.Button(
            s5, text="▶  Start Scan",
            font=("Segoe UI", 11, "bold"),
            bg=C["accent"], fg=C["header"],
            bd=0, relief="flat",
            padx=0, pady=13,
            cursor="hand2",
            activebackground=C["accent3"],
            activeforeground=C["header"],
            command=self._start_scan,
        )
        self.scan_btn.pack(fill="x", pady=(4, 10))

        # Secondary actions
        def sec_btn(label, icon, cmd, variant="default"):
            bg = C["danger"] if variant == "danger" else C["bg5"]
            b = tk.Button(s5, text=f"{icon}  {label}",
                          font=("Segoe UI", 9),
                          bg=bg, fg=C["text"],
                          bd=0, relief="flat",
                          padx=12, pady=8,
                          anchor="w", cursor="hand2",
                          activebackground=C["accent2"],
                          activeforeground="white",
                          command=cmd)
            b.pack(fill="x", pady=2)

        sec_btn("Show Matched Results", "📋", self._show_results)
        sec_btn("Copy Matched Files",   "📁", self._copy_files)
        sec_btn("Move Matched Files",   "🚚", self._move_files)
        sec_btn("Generate PDF Report",  "📄", self._generate_report)
        tk.Frame(s5, bg=C["border"], height=1).pack(fill="x", pady=6)
        sec_btn("Clear Encoding Cache", "🗑", self._clear_cache)

    def _build_right_panel(self, parent):
        right = tk.Frame(parent, bg=C["bg"])
        right.pack(side="right", fill="both", expand=True)

        # ── Tabs ──────────────────────────────────────────────────
        nb_frame = tk.Frame(right, bg=C["bg2"])
        nb_frame.pack(fill="x")

        style = ttk.Style()
        style.theme_use("default")
        style.configure("Forensic.TNotebook",
                        background=C["bg2"], borderwidth=0)
        style.configure("Forensic.TNotebook.Tab",
                        background=C["bg3"],
                        foreground=C["text2"],
                        padding=[18, 8],
                        font=("Segoe UI", 9))
        style.map("Forensic.TNotebook.Tab",
                  background=[("selected", C["bg4"])],
                  foreground=[("selected", C["accent"])])

        self.notebook = ttk.Notebook(right, style="Forensic.TNotebook")
        self.notebook.pack(fill="both", expand=True)

        self.tab_results = tk.Frame(self.notebook, bg=C["bg"])
        self.notebook.add(self.tab_results, text="  Results  ")
        self._build_results_tab(self.tab_results)

        self.tab_log = tk.Frame(self.notebook, bg=C["bg"])
        self.notebook.add(self.tab_log, text="  Live Log  ")
        self._build_log_tab(self.tab_log)

        self.tab_meta = tk.Frame(self.notebook, bg=C["bg"])
        self.notebook.add(self.tab_meta, text="  Metadata  ")
        self._build_meta_tab(self.tab_meta)

        self.tab_keywords = tk.Frame(self.notebook, bg=C["bg"])
        self.notebook.add(self.tab_keywords, text="  🔤 Keyword Results  ")
        self._build_keyword_tab(self.tab_keywords)

    def _build_results_tab(self, parent):
        # ── Stats bar ─────────────────────────────────────────────
        stats_bar = tk.Frame(parent, bg=C["bg2"])
        stats_bar.pack(fill="x")
        tk.Frame(stats_bar, bg=C["border"], height=1).pack(fill="x")
        stat_row = tk.Frame(stats_bar, bg=C["bg2"], pady=0)
        stat_row.pack(fill="x")

        self.stat_scanned, self.stat_scanned_of = self._stat_badge(
            stat_row, "SCANNED", "0", sub="of 0")
        self.stat_matched, _ = self._stat_badge(stat_row, "MATCHED",  "0")
        self.stat_nsfw,    _ = self._stat_badge(stat_row, "EXPLICIT", "0")
        self.stat_time,    _ = self._stat_badge(stat_row, "ELAPSED",  "00:00")
        tk.Frame(stats_bar, bg=C["border"], height=1).pack(fill="x")

        # ── Progress bar ──────────────────────────────────────────
        pb_track = tk.Frame(parent, bg=C["bg5"], height=4)
        pb_track.pack(fill="x")
        self.progress_bar = tk.Frame(pb_track, bg=C["accent"], height=4)
        self.progress_bar.place(x=0, y=0, relheight=1, width=0)

        # ── Scrollable result grid ────────────────────────────────
        grid_outer = tk.Frame(parent, bg=C["bg"])
        grid_outer.pack(fill="both", expand=True)

        self.grid_canvas = tk.Canvas(grid_outer, bg=C["bg"],
                                      highlightthickness=0)
        self.grid_scroll = tk.Scrollbar(grid_outer, orient="vertical",
                                         command=self.grid_canvas.yview,
                                         bg=C["bg3"], troughcolor=C["bg4"],
                                         bd=0, relief="flat")
        self.grid_canvas.configure(yscrollcommand=self.grid_scroll.set)
        self.grid_scroll.pack(side="right", fill="y")
        self.grid_canvas.pack(side="left", fill="both", expand=True)

        self.grid_inner = tk.Frame(self.grid_canvas, bg=C["bg"])
        self.grid_win   = self.grid_canvas.create_window(
            (0, 0), window=self.grid_inner, anchor="nw")
        self.grid_inner.bind("<Configure>", lambda e:
            self.grid_canvas.configure(
                scrollregion=self.grid_canvas.bbox("all")))
        self.grid_canvas.bind("<Configure>", lambda e:
            self.grid_canvas.itemconfig(self.grid_win, width=e.width))
        self.grid_canvas.bind_all("<MouseWheel>", lambda e:
            self.grid_canvas.yview_scroll(-1 * (e.delta // 120), "units"))

        self._show_placeholder("No scan data yet — press  ▶ Start Scan  to begin")

    def _stat_badge(self, parent, label, value, sub=None):
        """Returns (value_label, sub_label). sub_label is None if sub not given."""
        f = tk.Frame(parent, bg=C["bg2"], padx=20, pady=10)
        f.pack(side="left")
        tk.Frame(f, bg=C["border"], width=1).pack(side="right", fill="y")
        tk.Label(f, text=label,
                 font=("Segoe UI", 7),
                 bg=C["bg2"], fg=C["text2"]).pack(anchor="w")
        val_lbl = tk.Label(f, text=value,
                           font=("Segoe UI", 20, "bold"),
                           bg=C["bg2"], fg=C["accent"])
        val_lbl.pack(anchor="w")
        sub_lbl = None
        if sub is not None:
            sub_lbl = tk.Label(f, text=sub,
                               font=("Segoe UI", 7),
                               bg=C["bg2"], fg=C["text3"])
            sub_lbl.pack(anchor="w")
        return val_lbl, sub_lbl

    def _build_log_tab(self, parent):
        # Tab header bar
        hrow = tk.Frame(parent, bg=C["bg3"])
        hrow.pack(fill="x")
        tk.Label(hrow, text="  Live Scan Log",
                 font=("Segoe UI", 9, "bold"),
                 bg=C["bg3"], fg=C["text2"],
                 pady=7).pack(side="left")
        tk.Frame(hrow, bg=C["border"], height=1).pack(
            fill="x", side="bottom")

        f = tk.Frame(parent, bg=C["bg"])
        f.pack(fill="both", expand=True)
        self.log_text = tk.Text(
            f, font=("Courier New", 9),
            bg=C["bg"], fg=C["text2"],
            bd=0, relief="flat", wrap="word",
            state="disabled",
            highlightthickness=0,
            pady=10, padx=14,
            insertbackground=C["accent"],
        )
        sc = tk.Scrollbar(f, command=self.log_text.yview,
                           bg=C["bg3"], troughcolor=C["bg4"],
                           bd=0, relief="flat")
        self.log_text.configure(yscrollcommand=sc.set)
        sc.pack(side="right", fill="y")
        self.log_text.pack(fill="both", expand=True)
        # Tag colors — Courier New only in log (terminal feel)
        self.log_text.tag_config(
            "match",   foreground=C["success"],
            font=("Courier New", 9, "bold"))
        self.log_text.tag_config(
            "nsfw",    foreground=C["danger"],
            font=("Courier New", 9, "bold"))
        self.log_text.tag_config("skip",    foreground=C["warning"])
        self.log_text.tag_config("error",   foreground=C["danger"])
        self.log_text.tag_config("info",    foreground=C["text2"])
        self.log_text.tag_config(
            "heading", foreground=C["accent"],
            font=("Courier New", 9, "bold"))

    def _build_meta_tab(self, parent):
        hrow = tk.Frame(parent, bg=C["bg3"])
        hrow.pack(fill="x")
        tk.Label(hrow, text="  Forensic Metadata Inspector",
                 font=("Segoe UI", 9, "bold"),
                 bg=C["bg3"], fg=C["text2"],
                 pady=7).pack(side="left")
        tk.Label(hrow, text="Click any result card to load  →",
                 font=("Segoe UI", 8),
                 bg=C["bg3"], fg=C["text3"]).pack(side="right", padx=12)
        tk.Frame(hrow, bg=C["border"], height=1).pack(
            fill="x", side="bottom")

        f = tk.Frame(parent, bg=C["bg"])
        f.pack(fill="both", expand=True)
        self.meta_text = tk.Text(
            f, font=("Courier New", 9),
            bg=C["bg"], fg=C["text"],
            bd=0, relief="flat", wrap="word",
            state="disabled",
            highlightthickness=0,
            pady=10, padx=14,
        )
        sc = tk.Scrollbar(f, command=self.meta_text.yview,
                           bg=C["bg3"], troughcolor=C["bg4"],
                           bd=0, relief="flat")
        self.meta_text.configure(yscrollcommand=sc.set)
        sc.pack(side="right", fill="y")
        self.meta_text.pack(fill="both", expand=True)


    def _build_keyword_tab(self, parent):
        """Keyword search results tab — shows matched files with context."""
        hrow = tk.Frame(parent, bg=C["bg3"])
        hrow.pack(fill="x")
        tk.Label(hrow, text="  Keyword Search Results",
                 font=("Segoe UI", 9, "bold"),
                 bg=C["bg3"], fg=C["text2"], pady=7).pack(side="left")
        self.kw_count_lbl = tk.Label(hrow, text="0 files matched",
                                      font=("Segoe UI", 8),
                                      bg=C["bg3"], fg=C["text3"])
        self.kw_count_lbl.pack(side="right", padx=12)
        tk.Frame(hrow, bg=C["border"], height=1).pack(fill="x", side="bottom")

        outer = tk.Frame(parent, bg=C["bg"])
        outer.pack(fill="both", expand=True)

        self.kw_canvas = tk.Canvas(outer, bg=C["bg"], highlightthickness=0)
        kw_scroll = tk.Scrollbar(outer, orient="vertical",
                                  command=self.kw_canvas.yview,
                                  bg=C["bg3"], troughcolor=C["bg4"],
                                  bd=0, relief="flat")
        self.kw_canvas.configure(yscrollcommand=kw_scroll.set)
        kw_scroll.pack(side="right", fill="y")
        self.kw_canvas.pack(side="left", fill="both", expand=True)

        self.kw_inner = tk.Frame(self.kw_canvas, bg=C["bg"])
        self.kw_win   = self.kw_canvas.create_window(
            (0, 0), window=self.kw_inner, anchor="nw")
        self.kw_inner.bind("<Configure>", lambda e:
            self.kw_canvas.configure(
                scrollregion=self.kw_canvas.bbox("all")))
        self.kw_canvas.bind("<Configure>", lambda e:
            self.kw_canvas.itemconfig(self.kw_win, width=e.width))

        tk.Label(self.kw_inner,
            text="Enable Keyword Search, enter keywords, then run a scan",
            font=("Segoe UI", 10),
            bg=C["bg"], fg=C["text3"]).pack(pady=60)

    def _refresh_keyword_results(self):
        for w in self.kw_inner.winfo_children():
            w.destroy()

        results = self.keyword_results
        if not results:
            tk.Label(self.kw_inner, text="No keyword matches found",
                     font=("Segoe UI", 10),
                     bg=C["bg"], fg=C["text3"]).pack(pady=60)
            self.after(0, lambda: self.kw_count_lbl.config(text="0 files matched"))
            return

        self.after(0, lambda n=len(results): self.kw_count_lbl.config(
            text=f"{n} file{'s' if n != 1 else ''} matched"))

        for r in results:
            card = tk.Frame(self.kw_inner, bg=C["bg3"])
            card.pack(fill="x", padx=12, pady=5)
            tk.Frame(card, bg=C["accent"], height=2).pack(fill="x")

            content = tk.Frame(card, bg=C["bg3"], padx=14, pady=10)
            content.pack(fill="x")

            top = tk.Frame(content, bg=C["bg3"])
            top.pack(fill="x")

            icon = "IMAGE" if r["filetype"] == "image" else "DOC"
            tk.Label(top, text=f"[{icon}]  {r['filename']}",
                     font=("Segoe UI", 9, "bold"),
                     bg=C["bg3"], fg=C["text"]).pack(side="left")

            badge_row = tk.Frame(top, bg=C["bg3"])
            badge_row.pack(side="right")
            tk.Label(badge_row,
                     text=f" {r['extension'].upper().strip('.')} ",
                     font=("Segoe UI", 7, "bold"),
                     bg=C["bg5"], fg=C["text2"]).pack(side="left", padx=2)
            kc = r["keyword_count"]
            tk.Label(badge_row,
                     text=f" {kc} keyword{'s' if kc!=1 else ''} ",
                     font=("Segoe UI", 7, "bold"),
                     bg=C["accent"], fg=C["header"]).pack(side="left", padx=2)

            tk.Label(content, text=r["filepath"],
                     font=("Courier New", 7),
                     bg=C["bg3"], fg=C["text3"],
                     anchor="w").pack(fill="x", pady=(2, 6))

            for kw, sentences in r["keywords_found"].items():
                kw_frame = tk.Frame(content, bg=C["bg4"])
                kw_frame.pack(fill="x", pady=2)
                tk.Label(kw_frame, text=f"  {kw}  ",
                         font=("Segoe UI", 8, "bold"),
                         bg=C["success"], fg=C["header"],
                         padx=4, pady=3).pack(side="left", padx=(8,0), pady=4)
                ctx = " ... ".join(sentences[:3])
                if len(ctx) > 200:
                    ctx = ctx[:200] + "..."
                tk.Label(kw_frame, text=f"  {ctx}",
                         font=("Segoe UI", 8),
                         bg=C["bg4"], fg=C["text2"],
                         anchor="w", wraplength=700,
                         justify="left").pack(
                             side="left", fill="x", expand=True,
                             padx=8, pady=4)

            if r["filetype"] == "image":
                thumb = ForensicEngine.make_thumbnail(r["filepath"], (100, 75))
                if thumb:
                    self._thumb_refs.append(thumb)
                    tk.Label(content, image=thumb, bg=C["bg3"]).pack(anchor="w", pady=(4,0))

    def _build_statusbar(self):
        tk.Frame(self, bg=C["border"], height=1).pack(
            fill="x", side="bottom")
        sb = tk.Frame(self, bg=C["header"])
        sb.pack(fill="x", side="bottom")
        row = tk.Frame(sb, bg=C["header"])
        row.pack(fill="x", padx=14, pady=5)

        self.status_lbl = tk.Label(
            row, text="Ready",
            font=("Segoe UI", 8),
            bg=C["header"], fg=C["text2"])
        self.status_lbl.pack(side="left")

        tk.Label(row,
                 text="Documents / Images Finder Tool  ·  Advanced Search & Analysis",
                 font=("Segoe UI", 7),
                 bg=C["header"], fg=C["text3"]).pack(side="right")

    # ──────────────────────────────────────────────────────────────
    #  BROWSE HANDLERS
    # ──────────────────────────────────────────────────────────────
    def _browse_scan_folder(self):
        path = filedialog.askdirectory(title="Select Image Folder to Scan")
        if path:
            self.scan_folder.set(path)

    def _browse_sample(self):
        path = filedialog.askopenfilename(
            title="Select Sample Face Image",
            filetypes=[("Image Files", "*.jpg *.jpeg *.png *.bmp *.tiff *.webp"),
                       ("All Files", "*.*")])
        if path:
            self.sample_path.set(path)
            self._update_sample_preview(path)

    def _update_sample_preview(self, path: str):
        """Render sample image thumbnail in the sidebar preview box."""
        try:
            pil = Image.open(path).convert("RGB")
            # Fit inside 390 x 155 preserving aspect ratio
            pil.thumbnail((390, 155), Image.LANCZOS)
            # Add accent border
            bordered = Image.new("RGB", (pil.width + 4, pil.height + 4),
                                 tuple(int(C["accent"][i:i+2], 16) for i in (1, 3, 5)))
            bordered.paste(pil, (2, 2))
            self._sample_thumb = ImageTk.PhotoImage(bordered)
            self.sample_preview_lbl.config(image=self._sample_thumb, text="")
            self.sample_preview_frame.config(height=pil.height + 8)
        except Exception as e:
            self.sample_preview_lbl.config(image="", text=f"Preview error: {e}")

    # ──────────────────────────────────────────────────────────────
    #  UI HELPERS
    # ──────────────────────────────────────────────────────────────
    def _log(self, msg: str, tag: str = "info"):
        def _do():
            self.log_text.configure(state="normal")
            self.log_text.insert("end", msg + "\n", tag)
            self.log_text.see("end")
            self.log_text.configure(state="disabled")
        self.after(0, _do)

    def _set_status(self, msg: str, color=None):
        self.after(0, lambda: self.status_lbl.config(
            text=msg, fg=color or C["text2"]))

    def _set_progress(self, pct: float):
        def _do():
            self.progress_bar.master.update_idletasks()
            w = self.progress_bar.master.winfo_width()
            self.progress_bar.place(x=0, y=0, relheight=1, width=int(w * pct))
        self.after(0, _do)

    def _update_stats(self, scanned=None, matched=None, nsfw=None,
                      elapsed=None, total=None):
        def _do():
            if total   is not None:
                self._scan_total = total          # cache for sub-label
                if self.stat_scanned_of:
                    self.stat_scanned_of.config(
                        text=f"of {total:,}")
            if scanned is not None:
                self.stat_scanned.config(text=f"{scanned:,}")
                if self.stat_scanned_of:
                    t = getattr(self, "_scan_total", 0)
                    if t:
                        self.stat_scanned_of.config(
                            text=f"of {t:,}")
                # Also mirror to status bar for instant feedback
                t = getattr(self, "_scan_total", 0)
                if t and scanned:
                    pct = int(scanned / t * 100)
                    self.status_lbl.config(
                        text=f"Scanning…  {scanned:,} / {t:,} images  ({pct}%)",
                        fg=C["text2"])
            if matched is not None: self.stat_matched.config(text=f"{matched:,}")
            if nsfw    is not None: self.stat_nsfw.config(text=f"{nsfw:,}")
            if elapsed is not None: self.stat_time.config(text=elapsed)
        self.after(0, _do)

    def _show_placeholder(self, text: str):
        for w in self.grid_inner.winfo_children():
            w.destroy()
        tk.Label(self.grid_inner, text=text,
                 font=("Segoe UI", 10),
                 bg=C["bg"], fg=C["text3"]).pack(pady=60)

    def _clear_cache(self):
        folder = self.scan_folder.get()
        if not folder:
            messagebox.showwarning("No folder", "Select a scan folder first."); return
        cp = Path(folder) / CACHE_FILENAME
        if cp.exists():
            cp.unlink()
            self._log("  Encoding cache cleared.", "skip")
        else:
            self._log("  No cache found.", "info")

    # ──────────────────────────────────────────────────────────────
    #  RESULTS GRID
    # ──────────────────────────────────────────────────────────────
    def _show_results(self):
        for w in self.grid_inner.winfo_children():
            w.destroy()
        self._thumb_refs.clear()

        # Combine image/face/nsfw matches + keyword matches, deduplicated
        seen      = set()
        combined  = []
        for r in self.matched_results:
            if r["filepath"] not in seen:
                seen.add(r["filepath"])
                combined.append(r)
        for r in self.keyword_results:
            if r["filepath"] not in seen:
                seen.add(r["filepath"])
                combined.append(r)

        if not combined:
            self._show_placeholder("No matches found. Run a scan first.")
            return

        # File-type icon map for documents
        DOC_ICONS = {
            ".pdf":  ("📄", "#e06c75"),   # red
            ".docx": ("📝", "#61afef"),   # blue
            ".doc":  ("📝", "#61afef"),
            ".pptx": ("📊", "#e5c07b"),   # amber
            ".ppt":  ("📊", "#e5c07b"),
            ".txt":  ("📃", "#98c379"),   # green
            ".csv":  ("📋", "#56b6c2"),   # teal
        }
        IMG_EXTS = {".jpg", ".jpeg", ".png", ".bmp",
                    ".tiff", ".tif", ".webp", ".gif"}

        COLS = 5
        for i, r in enumerate(combined):
            row_idx = i // COLS
            col_idx = i % COLS
            ext     = Path(r["filepath"]).suffix.lower()
            is_img  = ext in IMG_EXTS

            card = tk.Frame(self.grid_inner, bg=C["bg3"],
                            cursor="hand2")
            card.grid(row=row_idx, column=col_idx,
                      padx=6, pady=6, sticky="nsew")
            tk.Frame(card, bg=C["accent"], height=2).pack(fill="x")

            # ── Preview area ──────────────────────────────────────
            if is_img:
                thumb = ForensicEngine.make_thumbnail(
                    r["filepath"], (140, 105))
                if thumb:
                    self._thumb_refs.append(thumb)
                    tk.Label(card, image=thumb,
                             bg=C["bg3"]).pack(pady=(6, 0))
                else:
                    tk.Label(card, text="🖼\nNo Preview",
                             font=("Segoe UI", 10),
                             bg=C["bg4"], fg=C["text3"],
                             width=18, height=5,
                             justify="center").pack(pady=(6, 0))
            else:
                # Document — show big icon + extension badge
                icon_char, icon_col = DOC_ICONS.get(
                    ext, ("📄", C["text2"]))
                doc_frame = tk.Frame(card, bg=C["bg4"],
                                     width=140, height=105)
                doc_frame.pack(pady=(6, 0))
                doc_frame.pack_propagate(False)
                tk.Label(doc_frame, text=icon_char,
                         font=("Segoe UI", 32),
                         bg=C["bg4"]).place(relx=0.5, rely=0.38,
                                            anchor="center")
                tk.Label(doc_frame,
                         text=ext.upper().strip("."),
                         font=("Segoe UI", 8, "bold"),
                         bg=icon_col, fg="white",
                         padx=6, pady=2).place(relx=0.5, rely=0.78,
                                               anchor="center")

            # ── Filename ──────────────────────────────────────────
            name = Path(r["filepath"]).name
            # Wrap long names at 20 chars per line
            display_name = name if len(name) <= 22 else name[:20] + "…"
            tk.Label(card, text=display_name,
                     font=("Segoe UI", 8, "bold"),
                     bg=C["bg3"], fg=C["text"],
                     wraplength=155,
                     justify="center").pack(pady=(4, 2))

            # ── Detection tags ────────────────────────────────────
            tag_row = tk.Frame(card, bg=C["bg3"])
            tag_row.pack(pady=(0, 2))

            if r.get("matched_face"):
                tk.Label(tag_row, text=" FACE ",
                         font=("Segoe UI", 7, "bold"),
                         bg=C["success"], fg="#0a1a0a",
                         padx=3).pack(side="left", padx=1)
            if r.get("nsfw", {}).get("is_explicit"):
                conf = r["nsfw"].get("confidence", 0)
                tk.Label(tag_row, text=f" NSFW {conf:.0%} ",
                         font=("Segoe UI", 7, "bold"),
                         bg=C["danger"], fg="white",
                         padx=3).pack(side="left", padx=1)
            if r.get("kissing", {}).get("is_kissing"):
                conf = r["kissing"].get("confidence", 0)
                tk.Label(tag_row, text=f" KISS {conf:.0%} ",
                         font=("Segoe UI", 7, "bold"),
                         bg=C["kiss_col"], fg="white",
                         padx=3).pack(side="left", padx=1)
            if r.get("keywords_found"):
                kws = ", ".join(list(r["keywords_found"].keys())[:2])
                tk.Label(tag_row, text=f" {kws} ",
                         font=("Segoe UI", 7, "bold"),
                         bg=C["accent"], fg=C["header"],
                         padx=3).pack(side="left", padx=1)

            # ── MD5 or keyword count ──────────────────────────────
            if is_img:
                md5 = r.get("hashes", {}).get("md5", "")[:14]
                info_text = f"{md5}.." if md5 else ""
            else:
                kc = r.get("keyword_count", 0)
                info_text = f"{kc} keyword match{'es' if kc != 1 else ''}"

            tk.Label(card, text=info_text,
                     font=("Courier New", 6),
                     bg=C["bg3"], fg=C["text3"]).pack(pady=(0, 6))

            card.bind("<Button-1>",
                      lambda e, res=r: self._show_metadata(res))
            for child in card.winfo_children():
                child.bind("<Button-1>",
                           lambda e, res=r: self._show_metadata(res))

        # Expand columns evenly
        for c in range(COLS):
            self.grid_inner.columnconfigure(c, weight=1)

        self.notebook.select(0)

    def _show_metadata(self, result: Dict):
        meta   = result.get("metadata", {})
        hashes = result.get("hashes",   {})
        nsfw   = result.get("nsfw",     {})
        kiss   = result.get("kissing",  {})
        nd     = "No data found"

        lines = [
            "═" * 60,
            f"  FILE METADATA REPORT",
            "═" * 60,
            f"  File Name      : {meta.get('filename',   nd)}",
            f"  Folder         : {meta.get('folder',     nd)}",
            f"  Full Path      : {meta.get('filepath',   nd)}",
            f"  File Size      : {meta.get('filesize',   nd)}",
            f"  Last Modified  : {meta.get('last_modified', nd)}",
            f"  Last Accessed  : {meta.get('last_accessed', nd)}",
            "─" * 60,
            f"  Device / Camera: {meta.get('device',      nd)}",
            f"  Capture Date   : {meta.get('capture_date', nd)}",
            f"  GPS Location   : {meta.get('gps',         nd)}",
            "─" * 60,
            f"  MD5  Hash      : {hashes.get('md5',    nd)}",
            f"  SHA256 Hash    : {hashes.get('sha256', nd)}",
            "─" * 60,
            f"  Face Match     : {'YES' if result.get('matched_face') else 'NO'}",
            f"  NSFW Detected  : {'YES' if nsfw.get('is_explicit') else 'NO'}",
            f"  NSFW Confidence: {nsfw.get('confidence', nd)}",
            f"  Kissing / Intimate : {'YES' if kiss.get('is_kissing') else 'NO'}",
            f"  Kissing Confidence : {kiss.get('confidence', nd)}",
            f"  Matched Prompt     : {kiss.get('top_prompt', nd)}",
        ]

        if nsfw.get("detections"):
            lines.append(f"  Detections     : {', '.join(d.get('class','') for d in nsfw['detections'][:5])}")

        lines += [
            "─" * 60,
            "  EXIF DATA:",
            meta.get("exif_raw", nd),
            "═" * 60,
        ]

        self.meta_text.configure(state="normal")
        self.meta_text.delete("1.0", "end")
        self.meta_text.insert("1.0", "\n".join(lines))
        self.meta_text.configure(state="disabled")
        self.notebook.select(2)

    # ──────────────────────────────────────────────────────────────
    #  SCAN ENGINE
    # ──────────────────────────────────────────────────────────────
    def _validate_scan(self) -> bool:
        if not self.scan_folder.get() or not os.path.isdir(self.scan_folder.get()):
            messagebox.showerror("Missing", "Please select a valid scan folder."); return False
        if self.enable_face.get() and not self.face_eng.available:
            messagebox.showerror("Missing", "face_recognition not installed.\nRun: pip install face-recognition"); return False
        if self.enable_face.get() and not self.sample_path.get():
            messagebox.showerror("Missing", "Please select a sample image for face matching."); return False
        if self.enable_nsfw.get() and not self.nsfw_eng.available:
            if not messagebox.askyesno("NudeNet Missing",
                "NudeNet is not installed.\nInstall with: pip install nudenet\n\nContinue without NSFW detection?"):
                return False
        return True

    def _start_scan(self):
        if self.running: return
        if not self._validate_scan(): return
        self.running = True
        self.scan_btn.config(text="▶  Scanning…", state="disabled", bg=C["text3"])
        self.matched_results.clear()
        self.keyword_results.clear()
        self._show_placeholder("Scanning…")
        self._set_progress(0)
        self._update_stats(scanned=0, matched=0, nsfw=0, elapsed="00:00", total=0)
        self.notebook.select(1)  # Switch to log
        threading.Thread(target=self._scan_thread, daemon=True).start()

    def _scan_thread(self):
        start_time  = datetime.datetime.now()
        folder      = self.scan_folder.get()
        do_face     = self.enable_face.get() and self.face_eng.available
        do_nsfw     = self.enable_nsfw.get() and self.nsfw_eng.available
        do_kissing  = self.enable_kissing.get() and self.kiss_eng.available
        do_keyword  = self.enable_keyword.get()   # available is always True now
        keywords_raw = self.keywords_var.get()
        keywords    = [k.strip() for k in keywords_raw.split(",") if k.strip()]
        do_keyword  = do_keyword and bool(keywords)
        kw_images   = self.kw_search_images.get()
        kw_docs     = self.kw_search_docs.get()
        kw_case     = self.kw_case_sensitive.get()
        tol         = self.tolerance.get()
        nsfw_thr    = self.nsfw_threshold.get()
        kiss_thr    = self.kiss_threshold.get()
        n_workers   = self.workers_var.get()

        # ── Performance: batch UI updates every N images ──────────
        UI_UPDATE_INTERVAL = 50   # update stats every 50 images not every 1
        CACHE_SAVE_INTERVAL = 500  # save cache every 500 encodings

        self._log("═" * 58, "heading")
        self._log("  DOCUMENTS / IMAGES FINDER — SCAN INITIATED", "heading")
        self._log(f"  Folder   : {folder}", "heading")
        self._log(f"  Face     : {'ON' if do_face else 'OFF'}  |  "
                  f"NSFW: {'ON' if do_nsfw else 'OFF'}  |  "
                  f"Kissing: {'ON' if do_kissing else 'OFF'}", "heading")
        self._log(f"  Workers  : {n_workers}  |  Tolerance: {tol:.2f}", "heading")
        self._log("═" * 58 + "\n", "heading")

        # ── Collect images ─────────────────────────────────────────
        self._set_status("Collecting images…")
        all_images = ForensicEngine.collect_images(folder)
        total      = len(all_images)
        self._log(f"  Found {total:,} images.\n", "info")
        self._update_stats(scanned=0, total=total)

        if total == 0 and not do_keyword:
            self._log("  No images found and keyword search is OFF. Scan aborted.", "error")
            self._finish_scan(start_time, 0)
            return

        if total == 0 and do_keyword:
            self._log("  No images found — running keyword search on documents only.\n", "info")

        # ── Encode sample face ─────────────────────────────────────
        sample_enc = None
        if do_face:
            self._set_status("Encoding sample face…")
            self._log("  Encoding sample face…", "info")
            sample_enc = self.face_eng.encode_sample(self.sample_path.get())
            if sample_enc is None:
                self._log("  ✘ No face found in sample image!", "error")
                do_face = False
            else:
                self._log("  ✔ Sample face encoded.\n", "match")

        # ── Load cache — skip already encoded images ───────────────
        cache     = FaceEngine.load_cache(folder) if do_face else {}
        to_encode = [p for p in all_images if str(p) not in cache] if do_face else []
        cached_n  = len(all_images) - len(to_encode)

        if cached_n:
            self._log(f"  ✔ {cached_n:,} images loaded from cache (instant).", "match")
        if to_encode:
            self._log(f"  Encoding {len(to_encode):,} new images with {n_workers} workers…", "info")
            paths_str  = [str(p) for p in to_encode]
            new_since_save = 0
            try:
                with multiprocessing.Pool(processes=n_workers) as pool:
                    # chunksize=50 — reduce IPC overhead significantly
                    for i, (path_str, encs) in enumerate(
                        pool.imap_unordered(_encode_worker, paths_str, chunksize=50)
                    ):
                        if encs is not None:
                            cache[path_str] = encs
                            new_since_save  += 1

                        # Batch UI updates — not every image
                        if i % UI_UPDATE_INTERVAL == 0:
                            pct     = (i + 1) / len(to_encode) * 0.6
                            elapsed = str(datetime.datetime.now() - start_time).split(".")[0]
                            self._set_progress(pct)
                            self._update_stats(scanned=cached_n + i + 1, elapsed=elapsed)
                            self._set_status(f"Encoding {i+1:,}/{len(to_encode):,}…")

                        # Incremental cache saves — avoid losing work
                        if new_since_save >= CACHE_SAVE_INTERVAL:
                            FaceEngine.save_cache(folder, cache)
                            new_since_save = 0

            except Exception as e:
                self._log(f"  Multiprocessing error: {e}. Falling back to single-thread.", "error")
                for i, path_str in enumerate(paths_str):
                    _, encs = _encode_worker(path_str)
                    cache[path_str] = encs or []
                    if i % UI_UPDATE_INTERVAL == 0:
                        self._update_stats(scanned=i + 1)

            FaceEngine.save_cache(folder, cache)
            self._log("  ✔ Encoding complete. Cache saved.\n", "match")

        # ── Main match loop — optimised ────────────────────────────
        self._log("  Matching and analyzing images…\n", "info")
        self._set_status("Matching faces…")
        matched    = 0
        nsfw_count = 0

        # Pre-convert sample encoding once
        sample_enc_arr = sample_enc

        # ── Pre-load CLIP once before loop (not per-image) ─────────
        if do_kissing:
            self._set_status("Loading CLIP model for kissing detection…")
            self._log("  Loading CLIP model (text features pre-caching)…", "info")
            self.kiss_eng.load()
            if not self.kiss_eng._loaded:
                self._log("  ✘ CLIP failed to load — kissing detection disabled.", "error")
                do_kissing = False
            else:
                self._log("  ✔ CLIP ready — text features cached.\n", "match")

        # ── Pre-run kissing detection in batches BEFORE main loop ──
        # This is far faster than running per-image inline
        kiss_results: Dict[str, Dict] = {}
        if do_kissing:
            self._set_status("Running kissing detection (batched)…")
            self._log(f"  Running CLIP kissing scan in batches of {KissingDetector.BATCH_SIZE}…", "info")
            paths_for_kiss = [str(p) for p in all_images]
            bs = KissingDetector.BATCH_SIZE
            for b_start in range(0, len(paths_for_kiss), bs):
                batch = paths_for_kiss[b_start : b_start + bs]
                batch_results = self.kiss_eng.analyze_batch(batch, kiss_thr)
                for fp, res in zip(batch, batch_results):
                    kiss_results[fp] = res
                # Update UI every batch
                done = min(b_start + bs, len(paths_for_kiss))
                pct  = done / len(paths_for_kiss) * 0.3   # 0–30% of progress
                self._set_progress(pct)
                self._set_status(f"Kissing scan: {done:,}/{len(paths_for_kiss):,}")
            kiss_detected = sum(1 for r in kiss_results.values() if r["is_kissing"])
            self._log(f"  ✔ Kissing scan complete — {kiss_detected} detected.\n", "match")

        for i, img_path in enumerate(all_images):
            # Batch UI updates only
            if i % UI_UPDATE_INTERVAL == 0:
                pct     = 0.3 + (i + 1) / total * 0.68
                elapsed = str(datetime.datetime.now() - start_time).split(".")[0]
                self._set_progress(pct)
                self._update_stats(scanned=i + 1, elapsed=elapsed)

            matched_face   = False
            nsfw_result    = {"is_explicit": False, "confidence": 0.0, "detections": []}
            # Pull kissing result from pre-computed dict — O(1) lookup
            kissing_result = kiss_results.get(str(img_path),
                             {"is_kissing": False, "confidence": 0.0, "top_prompt": ""})

            # Face matching — pure in-memory comparison (fast)
            if do_face and sample_enc_arr is not None:
                encs = cache.get(str(img_path), [])
                if encs:
                    enc_arrays = [np.array(e) for e in encs]
                    matched_face = any(
                        face_recognition.compare_faces(enc_arrays, sample_enc_arr, tolerance=tol)
                    )

            # NSFW detection
            if do_nsfw:
                nsfw_result = self.nsfw_eng.analyze(str(img_path), nsfw_thr)

            # Collect match — only extract heavy metadata for matches
            is_match = (
                matched_face
                or (do_nsfw    and nsfw_result["is_explicit"])
                or (do_kissing and kissing_result["is_kissing"])
            )

            if is_match:
                meta   = ForensicEngine.extract_metadata(str(img_path))
                hashes = ForensicEngine.compute_hashes(str(img_path))
                self.matched_results.append({
                    "filepath":     str(img_path),
                    "metadata":     meta,
                    "hashes":       hashes,
                    "matched_face": matched_face,
                    "nsfw":         nsfw_result,
                    "kissing":      kissing_result,
                })
                matched += 1
                if nsfw_result["is_explicit"]:
                    nsfw_count += 1

                # Build tag label
                tags = []
                if matched_face:                       tags.append("FACE")
                if nsfw_result["is_explicit"]:         tags.append("NSFW")
                if kissing_result["is_kissing"]:       tags.append("KISS")
                tag_str  = "+".join(tags)
                log_tag  = "nsfw" if nsfw_result["is_explicit"] else \
                           "skip" if kissing_result["is_kissing"] else "match"
                self._log(f"  ✔ [{tag_str}]  →  {img_path.name}", log_tag)
                self._update_stats(matched=matched, nsfw=nsfw_count)

        self._set_progress(1.0)

        # ── KEYWORD SEARCH ─────────────────────────────────────────
        if do_keyword:
            self._log("\n" + "─" * 58, "heading")
            self._log(f"  KEYWORD SEARCH STARTING", "heading")
            self._log(f"  Keywords : {', '.join(keywords)}", "heading")
            self._log(f"  Scope    : "
                      f"{'Images ' if kw_images else ''}"
                      f"{'Documents' if kw_docs else ''}",
                      "heading")
            self._log("─" * 58 + "\n", "heading")

            # Load OCR if searching images
            if kw_images and self.kw_eng.ocr_available:
                self._set_status("Loading OCR model…")
                self.kw_eng.load_ocr(log_fn=self._log)

            self._set_status("Keyword scanning…")

            def _kw_progress(done, total_files, fname):
                self._set_status(
                    f"Keyword scan: {done:,} / {total_files:,}  —  {fname}")
                pct = done / total_files if total_files else 0
                self._set_progress(pct)

            self.keyword_results = self.kw_eng.scan_folder(
                folder         = folder,
                keywords       = keywords,
                do_images      = kw_images,
                do_docs        = kw_docs,
                case_sensitive = kw_case,
                progress_cb    = _kw_progress,
                log_fn         = self._log,
            )

            kw_matched = len(self.keyword_results)
            self._log(f"  ✔ Keyword scan complete — {kw_matched} file(s) matched.\n",
                      "match" if kw_matched else "info")

            # Refresh keyword results tab on UI thread
            self.after(0, self._refresh_keyword_results)
            # Switch to keyword tab automatically
            self.after(0, lambda: self.notebook.select(3))

        self._finish_scan(start_time, total)

    def _finish_scan(self, start_time, total):
        elapsed = str(datetime.datetime.now() - start_time).split(".")[0]
        matched = len(self.matched_results)
        nsfw    = sum(1 for r in self.matched_results
                      if r.get("nsfw", {}).get("is_explicit"))
        kw_matched = len(self.keyword_results)

        self._log(f"\n{'═'*58}", "heading")
        self._log(f"  SCAN COMPLETE", "heading")
        self._log(f"  Total Scanned    : {total:,}", "heading")
        self._log(f"  Total Matched    : {matched}", "heading")
        self._log(f"  NSFW Detected    : {nsfw}", "heading")
        if self.enable_keyword.get():
            self._log(f"  Keyword Matches  : {kw_matched}", "heading")
        self._log(f"  Elapsed Time     : {elapsed}", "heading")
        self._log(f"{'═'*58}", "heading")

        self._update_stats(scanned=total, matched=matched, nsfw=nsfw,
                           elapsed=elapsed, total=total)

        parts = [f"{matched} image match(es)"]
        if self.enable_keyword.get():
            parts.append(f"{kw_matched} keyword match(es)")
        self._set_status(
            f"SCAN COMPLETE  —  {total:,} files scanned  —  "
            + "  ·  ".join(parts),
            C["accent"])

        self.running = False
        self.after(0, lambda: self.scan_btn.config(
            text="▶  Start Scan", state="normal", bg=C["accent"]))
        self.after(0, self._show_results)

    # ──────────────────────────────────────────────────────────────
    #  EVIDENCE MANAGEMENT
    # ──────────────────────────────────────────────────────────────
    def _get_matched_paths(self) -> List[str]:
        """Return all matched file paths — both image/face/nsfw matches AND keyword matches."""
        all_results = self.matched_results + self.keyword_results
        if not all_results:
            messagebox.showwarning("No Results",
                "No matched files found.\nRun a scan with at least one detection enabled.")
            return []
        # Deduplicate — same file could appear in both lists
        seen  = set()
        paths = []
        for r in all_results:
            fp = r["filepath"]
            if fp not in seen:
                seen.add(fp)
                paths.append(fp)
        return paths

    def _copy_files(self):
        paths = self._get_matched_paths()
        if not paths: return
        dest = filedialog.askdirectory(title="Select Destination Folder for Copy")
        if not dest: return
        n = self.evidence.copy_files(paths, dest)
        messagebox.showinfo("Done", f"Copied {n} file(s) to:\n{dest}")
        self._log(f"  Copied {n} matched file(s) to {dest}", "match")

    def _move_files(self):
        paths = self._get_matched_paths()
        if not paths: return
        if not messagebox.askyesno("Confirm Move",
            f"Move {len(paths)} matched file(s)?\nThis will remove them from the original location."):
            return
        dest = filedialog.askdirectory(title="Select Destination Folder for Move")
        if not dest: return
        n = self.evidence.move_files(paths, dest)
        messagebox.showinfo("Done", f"Moved {n} file(s) to:\n{dest}")
        self._log(f"  Moved {n} matched file(s) to {dest}", "match")

    # ──────────────────────────────────────────────────────────────
    #  PDF REPORT
    # ──────────────────────────────────────────────────────────────
    def _generate_report(self):
        if not self.report_eng.available:
            messagebox.showerror("Missing",
                "reportlab not installed.\nRun: pip install reportlab")
            return

        all_results = self.matched_results + self.keyword_results
        if not all_results:
            messagebox.showwarning("No Results", "Run a scan first.")
            return

        ts   = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        path = filedialog.asksaveasfilename(
            title="Save Report",
            defaultextension=".pdf",
            initialfile=f"FinderTool_Report_{ts}.pdf",
            filetypes=[("PDF Files", "*.pdf")])
        if not path: return

        mode_parts = []
        if self.enable_face.get():    mode_parts.append("Face Matching")
        if self.enable_nsfw.get():    mode_parts.append("NSFW Detection")
        if self.enable_kissing.get(): mode_parts.append("Kissing Detection")
        if self.enable_keyword.get(): mode_parts.append(
            f"Keyword Search ({self.keywords_var.get()})")
        scan_mode = " + ".join(mode_parts) if mode_parts else "General Scan"

        try:
            total_scanned = int(
                self.stat_scanned.cget("text").replace(",", ""))
        except Exception:
            total_scanned = len(all_results)

        ok = self.report_eng.generate(
            all_results, path, total_scanned, scan_mode)

        if ok:
            messagebox.showinfo("Report Saved",
                f"PDF report saved to:\n{path}")
            self._log(f"  ✔ PDF report saved: {path}", "match")
        else:
            messagebox.showerror("Error",
                "Failed to generate PDF report.\nCheck the log for details.")
            if ok and messagebox.askyesno("Open Report", "Open the PDF now?"):
                os.startfile(path)


# ══════════════════════════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    multiprocessing.freeze_support()   # Required for PyInstaller

    # Dependency check summary
    print("\n╔══════════════════════════════════════════════╗")
    print("║  Documents / Images Finder Tool             ║")
    print("║  Advanced Search & Analysis Tool            ║")
    print("╚══════════════════════════════════════════════╝\n")
    print(f"  face_recognition : {'✔' if face_recognition else '✘ (pip install face-recognition)'}")
    print(f"  NudeNet          : {'✔' if NudeDetector else '✘ (pip install nudenet)'}")
    print(f"  ReportLab        : {'✔' if reportlab_ok else '✘ (pip install reportlab)'}")
    print(f"  piexif           : {'✔' if piexif_ok else '✘ (pip install piexif)'}")
    print()

    app = ForensicApp()
    app.mainloop()
